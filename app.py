# app.py
import io
import base64
import logging
import copy
from difflib import SequenceMatcher
from dataclasses import dataclass
from datetime import date, timedelta
from pathlib import Path
import re
import json
from pathlib import Path

import pandas as pd
from dash import (
    Dash,
    html,
    dcc,
    Input,
    Output,
    State,
    callback,
    no_update,
    ALL,
    callback_context,
)
from openpyxl import load_workbook
from openpyxl.utils import range_boundaries, get_column_letter
import numbers
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
from lxml import etree

app = Dash(__name__)
app.title = "Deck Automator (MVP)"
server = app.server
logger = logging.getLogger(__name__)
TEMPLATE_DIR = Path(__file__).resolve().parent / "assets" / "templates"
PROJECT_TEMPLATES = {
    "PnP": TEMPLATE_DIR / "PnP.pptx",
    "MMx": TEMPLATE_DIR / "MMx.pptx",
    "MMM": TEMPLATE_DIR / "MMM.pptx",
}
DISPLAY_LABEL = {"Own": "Own", "Cross": "Competitor"}
_NUM_CACHE_WARNING_LIMIT = 10
_num_cache_warning_count = 0


@dataclass(frozen=True)
class CompanyWeekMapper:
    """
    Maps between a sequential internal 'company_week' integer and YEARWK (YYYYWW),
    assuming YEARWK is an ISO week (ISO-8601).

    Provide at least one anchor mapping:
        (anchor_company_week -> anchor_yearwk)

    Optionally provide a second anchor to validate the mapping.
    """
    anchor_company_week: int
    anchor_yearwk: int
    check_company_week: int | None = None
    check_yearwk: int | None = None

    @staticmethod
    def _yearwk_to_monday(yearwk: int) -> date:
        y, w = divmod(yearwk, 100)
        if not (1 <= w <= 53):
            raise ValueError(f"Invalid YEARWK week number: {yearwk}")
        # Monday of ISO week
        return date.fromisocalendar(y, w, 1)

    @staticmethod
    def _monday_to_yearwk(d: date) -> int:
        iso_y, iso_w, _ = d.isocalendar()
        return iso_y * 100 + iso_w

    def __post_init__(self):
        # Optional consistency check with second anchor
        if (self.check_company_week is None) ^ (self.check_yearwk is None):
            raise ValueError("Provide both check_company_week and check_yearwk, or neither.")

        if self.check_company_week is not None:
            a_date = self._yearwk_to_monday(self.anchor_yearwk)
            delta = self.check_company_week - self.anchor_company_week
            derived = self._monday_to_yearwk(a_date + timedelta(weeks=delta))
            if derived != self.check_yearwk:
                raise ValueError(
                    f"Anchors inconsistent: derived {derived} but expected {self.check_yearwk}."
                )

    def to_yearwk(self, company_week: int) -> int:
        a_date = self._yearwk_to_monday(self.anchor_yearwk)
        delta_weeks = company_week - self.anchor_company_week
        out_date = a_date + timedelta(weeks=delta_weeks)
        return self._monday_to_yearwk(out_date)

    def to_company_week(self, yearwk: int) -> int:
        a_date = self._yearwk_to_monday(self.anchor_yearwk)
        target_date = self._yearwk_to_monday(yearwk)
        delta_days = (target_date - a_date).days

        if delta_days % 7 != 0:
            # Should never happen because both are Mondays, but keep it safe
            raise ValueError("Non-week-aligned difference; check inputs.")
        delta_weeks = delta_days // 7
        return self.anchor_company_week + delta_weeks

def bytes_from_contents(contents):
    _, content_string = contents.split(',')
    return base64.b64decode(content_string)

def _json_safe(x):
    try:
        import numpy as np
        if isinstance(x, (np.integer, np.floating)):
            return float(x)
    except Exception:
        pass
    return x


def df_from_contents(contents, filename):
    decoded = bytes_from_contents(contents)
    if filename.lower().endswith((".xlsx", ".xls", ".xlsb")):
        read_options = {}
        if filename.lower().endswith(".xlsb"):
            read_options["engine"] = "pyxlsb"
        return pd.read_excel(io.BytesIO(decoded), **read_options)
    elif filename.lower().endswith(".csv"):
        return pd.read_csv(io.StringIO(decoded.decode('utf-8')))
    else:
        raise ValueError("Unsupported file format. Please upload CSV or Excel.")


def raw_df_from_contents(contents, filename):
    decoded = bytes_from_contents(contents)
    if filename.lower().endswith((".xlsx", ".xls", ".xlsb")):
        read_options = {}
        if filename.lower().endswith(".xlsb"):
            read_options["engine"] = "pyxlsb"
        return pd.read_excel(io.BytesIO(decoded), header=None, **read_options)
    elif filename.lower().endswith(".csv"):
        return pd.read_csv(io.StringIO(decoded.decode("utf-8")), header=None)
    else:
        raise ValueError("Unsupported file format. Please upload CSV or Excel.")


def scope_df_from_contents(contents, filename):
    if not filename or not filename.lower().endswith((".xlsx", ".xlsb")):
        raise ValueError("Scope file must be an Excel workbook (.xlsx or .xlsb).")

    decoded = bytes_from_contents(contents)
    read_options = {}
    if filename.lower().endswith(".xlsb"):
        read_options["engine"] = "pyxlsb"
    scope_df = pd.read_excel(
        io.BytesIO(decoded),
        sheet_name="Overall Information",
        **read_options,
    )
    if scope_df.empty:
        return None
    return scope_df


def project_details_df_from_contents(contents, filename):
    if not filename or not filename.lower().endswith((".xlsx", ".xlsb")):
        raise ValueError("Scope file must be an Excel workbook (.xlsx or .xlsb).")

    decoded = bytes_from_contents(contents)
    read_options = {}
    if filename.lower().endswith(".xlsb"):
        read_options["engine"] = "pyxlsb"
    try:
        return pd.read_excel(
            io.BytesIO(decoded),
            sheet_name="Project Details",
            header=None,
            **read_options,
        )
    except ValueError:
        return None


def _normalize_project_details_label(text: object) -> str:
    if pd.isna(text):
        return ""
    normalized = str(text).strip().lower()
    normalized = " ".join(normalized.split())
    normalized = normalized.rstrip(" :;,.!?")
    return normalized


def _project_detail_value_from_df(
    project_details_df: pd.DataFrame | None,
    label_key: str,
    synonyms: list[str],
    canonical: str,
):
    if project_details_df is None or project_details_df.empty or project_details_df.shape[1] < 2:
        return None

    synonym_matches = []
    normalized_synonyms = {_normalize_project_details_label(item) for item in synonyms}
    for row_idx, row in project_details_df.iterrows():
        cell_value = row.iloc[0]
        normalized_cell = _normalize_project_details_label(cell_value)
        if not normalized_cell:
            continue
        if normalized_cell in normalized_synonyms:
            synonym_matches.append((row_idx, cell_value, normalized_cell, 1.0))

    if synonym_matches:
        candidates = synonym_matches
    else:
        candidates = []
        canonical_norm = _normalize_project_details_label(canonical)
        for row_idx, row in project_details_df.iterrows():
            cell_value = row.iloc[0]
            normalized_cell = _normalize_project_details_label(cell_value)
            if not normalized_cell:
                continue
            score = SequenceMatcher(None, normalized_cell, canonical_norm).ratio()
            if score >= 0.85:
                candidates.append((row_idx, cell_value, normalized_cell, score))

    if not candidates:
        raise ValueError(f"Could not find Project Details label for {label_key}")

    candidates.sort(key=lambda item: item[3], reverse=True)
    if len(candidates) > 1 and candidates[1][3] >= candidates[0][3] - 0.03:
        details = ", ".join(
            f"{str(item[1]).strip()} ({item[3]:.2f})" for item in candidates
        )
        raise ValueError(
            f"Ambiguous Project Details label for {label_key}. Candidates: {details}"
        )

    row_idx, original_label, _, _ = candidates[0]
    logger.info(
        "Matched Project Details label for %s: %s", label_key, str(original_label)
    )
    raw_value = project_details_df.iloc[row_idx, 1]
    if pd.isna(raw_value):
        return ""
    return str(raw_value).strip()


def product_description_df_from_contents(contents, filename):
    if not filename or not filename.lower().endswith((".xlsx", ".xlsb")):
        raise ValueError("Scope file must be an Excel workbook (.xlsx or .xlsb).")

    decoded = bytes_from_contents(contents)
    read_options = {}
    if filename.lower().endswith(".xlsb"):
        read_options["engine"] = "pyxlsb"
    with pd.ExcelFile(io.BytesIO(decoded), **read_options) as excel_file:
        target_sheet = _find_sheet_by_candidates(
            excel_file.sheet_names, "PRODUCT DESCRIPTION"
        )
        if not target_sheet:
            return None
        product_df = excel_file.parse(target_sheet)
    if product_df.empty:
        return None
    return product_df


def target_brand_from_scope_df(scope_df):
    if scope_df is None or scope_df.empty:
        return None

    column_lookup = {str(col).strip().lower(): col for col in scope_df.columns}
    if "target brand" in column_lookup:
        series = scope_df[column_lookup["target brand"]].dropna()
        if not series.empty:
            return str(series.iloc[0])

    for _, row in scope_df.iterrows():
        if not len(row):
            continue
        label = str(row.iloc[0]).strip().lower()
        normalized_label = label.rstrip(":")
        if normalized_label == "target brand" and len(row) > 1 and pd.notna(row.iloc[1]):
            return str(row.iloc[1])

    return None

def modelled_category_from_scope_df(scope_df):
    if scope_df is None or scope_df.empty:
        return None

    if scope_df.shape[1] >= 2:
        for _, row in scope_df.iterrows():
            if not len(row):
                continue
            label = str(row.iloc[0]).strip()
            normalized_label = _normalize_label(label)
            if normalized_label == "category" and pd.notna(row.iloc[1]):
                return str(row.iloc[1])

    return None


def _normalize_column_name(value: str) -> str:
    return "".join(ch for ch in value.strip().lower() if ch.isalnum())


def _find_sheet_by_candidates(sheet_names: list[str], target: str) -> str | None:
    normalized_target = _normalize_column_name(target)
    normalized_sheets = {
        _normalize_column_name(str(sheet_name)): sheet_name
        for sheet_name in sheet_names
    }
    if normalized_target in normalized_sheets:
        return normalized_sheets[normalized_target]
    for normalized_sheet, sheet_name in normalized_sheets.items():
        if normalized_target in normalized_sheet or normalized_sheet in normalized_target:
            return sheet_name
    from difflib import get_close_matches

    matches = get_close_matches(
        normalized_target,
        list(normalized_sheets.keys()),
        n=1,
        cutoff=0.7,
    )
    if matches:
        return normalized_sheets[matches[0]]
    return None


def _find_column_by_candidates(df: pd.DataFrame, candidates: list[str]):
    normalized_columns = {_normalize_column_name(str(col)): col for col in df.columns}
    candidate_normalized = [_normalize_column_name(candidate) for candidate in candidates]
    for candidate in candidate_normalized:
        if candidate in normalized_columns:
            return normalized_columns[candidate]
    for column_key, column_name in normalized_columns.items():
        for candidate in candidate_normalized:
            if candidate in column_key or column_key in candidate:
                return column_name
    from difflib import get_close_matches

    matches = get_close_matches(
        " ".join(candidate_normalized),
        list(normalized_columns.keys()),
        n=1,
        cutoff=0.75,
    )
    if matches:
        return normalized_columns[matches[0]]
    return None


def _find_column_by_row_values(row: pd.Series, candidates: list[str]):
    normalized_values = {}
    for column, value in row.items():
        if pd.isna(value):
            continue
        normalized_values[_normalize_column_name(str(value))] = column
    if not normalized_values:
        return None

    candidate_normalized = [_normalize_column_name(candidate) for candidate in candidates]
    for candidate in candidate_normalized:
        if candidate in normalized_values:
            return normalized_values[candidate]
    for value_key, column_name in normalized_values.items():
        for candidate in candidate_normalized:
            if candidate in value_key or value_key in candidate:
                return column_name
    from difflib import get_close_matches

    matches = get_close_matches(
        " ".join(candidate_normalized),
        list(normalized_values.keys()),
        n=1,
        cutoff=0.75,
    )
    if matches:
        return normalized_values[matches[0]]
    return None


def _is_target_flag(value):
    if pd.isna(value):
        return False
    try:
        return float(value) == 1
    except (TypeError, ValueError):
        text = str(value).strip().lower()
        return text in {"1", "yes", "y", "true", "t"}


def _target_values_from_scope(
    scope_df: pd.DataFrame,
    target_col_candidates: list[str],
    value_col_candidates: list[str],
):
    if scope_df is None or scope_df.empty:
        return None
    target_col = _find_column_by_candidates(scope_df, target_col_candidates)
    value_col = _find_column_by_candidates(scope_df, value_col_candidates)
    if not target_col or not value_col:
        return None

    values = []
    seen = set()
    for _, row in scope_df.iterrows():
        if not _is_target_flag(row[target_col]):
            continue
        value = row[value_col]
        if pd.isna(value):
            continue
        name = str(value).strip()
        if not name or name in seen:
            continue
        seen.add(name)
        values.append(name)
    return values or None


def target_brands_from_scope_df(scope_df: pd.DataFrame):
    return _target_values_from_scope(
        scope_df,
        ["Target Brand", "Target_Brand"],
        ["Brand", "Brand Name"],
    )


def target_manufacturers_from_scope_df(scope_df: pd.DataFrame):
    return _target_values_from_scope(
        scope_df,
        ["Target Manufacturer", "Target_Manufacturer", "Target Mfr", "Target Mfg"],
        ["Manufacturer", "Mfr", "Mfg"],
    )


def target_brands_from_product_description(product_df: pd.DataFrame):
    if product_df is None or product_df.empty:
        return None
    target_col = _find_column_by_candidates(product_df, ["Target Brand", "Target_Brand"])
    brand_col = _find_column_by_candidates(product_df, ["Brand"])
    if not target_col or not brand_col:
        return None

    brands = []
    seen = set()
    for _, row in product_df.iterrows():
        if not _is_target_flag(row[target_col]):
            continue
        brand_value = row[brand_col]
        if pd.isna(brand_value):
            continue
        brand_name = str(brand_value).strip()
        if not brand_name or brand_name in seen:
            continue
        seen.add(brand_name)
        brands.append(brand_name)
    return brands or None

def target_dimensions_from_product_description(product_df: pd.DataFrame) -> list[str]:
    if product_df is None or product_df.empty:
        return []

    normalized_columns = {
        _normalize_column_name(str(col)): col for col in product_df.columns
    }
    lines = []
    seen_dimensions = set()
    for column in product_df.columns:
        column_name = str(column)
        if not column_name.strip().lower().startswith("target"):
            continue
        base_name = column_name.strip()[len("target"):].lstrip(" _-").strip()
        if not base_name:
            continue
        base_key = _normalize_column_name(base_name)
        if base_key in seen_dimensions:
            continue
        base_column = normalized_columns.get(base_key) or _find_column_by_candidates(
            product_df, [base_name]
        )
        if not base_column:
            continue

        values = []
        seen_values = set()
        for _, row in product_df.iterrows():
            if not _is_target_flag(row[column]):
                continue
            value = row[base_column]
            if pd.isna(value):
                continue
            value_name = str(value).strip()
            if not value_name or value_name in seen_values:
                continue
            seen_values.add(value_name)
            values.append(value_name)

        if values:
            base_label = str(base_column).strip()
            lines.append(f"Target {base_label}(s): {', '.join(values)}")
            seen_dimensions.add(base_key)

    return lines

def target_lines_from_product_description(product_df: pd.DataFrame) -> list[str]:
    return target_dimensions_from_product_description(product_df)

def _find_slide_by_marker(prs, marker_text: str):
    marker_text = marker_text.strip()
    for slide in prs.slides:
        for shape in slide.shapes:
            shape_name = getattr(shape, "name", "") or ""
            if marker_text and marker_text in shape_name:
                return slide
            if shape.has_text_frame and marker_text in shape.text_frame.text:
                return slide
    return None


def _build_category_waterfall_df(
    gathered_df: pd.DataFrame,
    target_level_label: str | None = None,
) -> pd.DataFrame:
    header_row = gathered_df.iloc[0] if len(gathered_df) else None
    data_start_idx = 0
    vars_col, vars_idx, _ = _resolve_column_from_candidates(
        gathered_df,
        header_row,
        ["Vars", "Var", "Variable", "Variable Name", "Bucket", "Driver"],
        context="Vars/Variable",
    )
    data_start_idx = max(data_start_idx, vars_idx)
    if not vars_col:
        raise ValueError("The gatheredCN10 file is missing a Vars/Variable column for the waterfall.")

    series_candidates = {
        "Base": ["Base"],
        "Promo": ["Promo", "Promotion", "Promotions"],
        "Media": ["Media"],
        "Blanks": ["Blanks", "Blank"],
        "Positives": ["Positives", "Positive", "Pos"],
        "Negatives": ["Negatives", "Negative", "Neg"],
    }
    series_columns = {}
    missing = []
    for key, candidates in series_candidates.items():
        found, found_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            candidates,
            context=key,
        )
        if not found:
            missing.append(key)
        else:
            series_columns[key] = found
            data_start_idx = max(data_start_idx, found_idx)
    if missing:
        raise ValueError(
            "The gatheredCN10 file is missing waterfall columns: "
            + ", ".join(missing)
        )

    label_candidates = {
        "labs-Base": ["labs-Base", "labs Base", "labels-Base", "labels Base"],
        "labs-Promo": ["labs-Promo", "labs Promo", "labels-Promo", "labels Promo"],
        "labs-Media": ["labs-Media", "labs Media", "labels-Media", "labels Media"],
        "labs-Blanks": ["labs-Blanks", "labs Blanks", "labels-Blanks", "labels Blanks"],
        "labs-Positives": [
            "labs-Positives",
            "labs Positives",
            "labels-Positives",
            "labels Positives",
        ],
        "labs-Negatives": [
            "labs-Negatives",
            "labs Negatives",
            "labels-Negatives",
            "labels Negatives",
        ],
    }
    label_columns = {}
    for key, candidates in label_candidates.items():
        found, found_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            candidates,
            context=key,
        )
        if found:
            label_columns[key] = found
            data_start_idx = max(data_start_idx, found_idx)

    data_df = gathered_df.iloc[data_start_idx:].copy()
    if target_level_label:
        target_col, target_idx, _ = _resolve_column_from_candidates(
            gathered_df,
            header_row,
            ["Target Level Label", "Target Level"],
            context="Target Level Label",
        )
        if target_col:
            data_start_idx = max(data_start_idx, target_idx)
            data_df = gathered_df.iloc[data_start_idx:].copy()
            normalized_target = _normalize_text_value(target_level_label)
            target_series = data_df[target_col].map(_normalize_text_value)
            data_df = data_df[target_series == normalized_target]

    ordered_cols = [vars_col] + [series_columns[key] for key in series_candidates]
    ordered_cols += list(label_columns.values())
    waterfall_df = data_df.loc[:, ordered_cols].copy()
    rename_map = {vars_col: "Vars", **{v: k for k, v in series_columns.items()}}
    rename_map.update({v: k for k, v in label_columns.items()})
    waterfall_df = waterfall_df.rename(columns=rename_map)

    for key in series_candidates:
        waterfall_df[key] = pd.to_numeric(waterfall_df[key], errors="coerce").fillna(0)
    if "Negatives" in waterfall_df.columns:
        waterfall_df["Negatives"] = waterfall_df["Negatives"].abs()

    return waterfall_df

def _target_level_labels_from_gathered_df(gathered_df: pd.DataFrame) -> list[str]:
    if gathered_df is None or gathered_df.empty:
        return []
    label_col, data_start_idx = _target_level_label_column_exact(gathered_df)
    if not label_col:
        raise ValueError("The gatheredCN10 file is missing the Target Level Label column.")
    labels = (
        gathered_df.iloc[data_start_idx:][label_col]
        .dropna()
        .astype(str)
        .map(str.strip)
    )
    unique_labels = []
    seen = set()
    for label in labels:
        if not label or label in seen:
            continue
        seen.add(label)
        unique_labels.append(label)
    return unique_labels


def _target_level_label_column_exact(gathered_df: pd.DataFrame) -> tuple[str | None, int]:
    if gathered_df is None or gathered_df.empty:
        return None, 0
    column = _find_column_by_candidates(
        gathered_df,
        ["Target Level Label", "Target Level"],
    )
    if column:
        return column, 0
    header_row = gathered_df.iloc[0]
    column = _find_column_by_row_values(header_row, ["Target Level Label", "Target Level"])
    if column:
        return column, 1
    return None, 0


def _target_level_labels_from_gathered_df_with_filters(
    gathered_df: pd.DataFrame,
    year1: str | None = None,
    year2: str | None = None,
    target_labels: list[str] | None = None,
) -> list[str]:
    if gathered_df is None or gathered_df.empty:
        return []
    label_col, data_start_idx = _target_level_label_column_exact(gathered_df)
    if not label_col:
        raise ValueError("The gatheredCN10 file is missing the Target Level Label column.")
    data_df = gathered_df.iloc[data_start_idx:]
    filtered_df = data_df

    year_col = _find_column_by_candidates(gathered_df, ["Year", "Model Year"])
    if year_col and (year1 is not None or year2 is not None):
        normalized_years = {
            _normalize_text_value(value)
            for value in (year1, year2)
            if value is not None
        }
        if normalized_years:
            year_series = data_df[year_col].map(_normalize_text_value)
            filtered_df = filtered_df[year_series.isin(normalized_years)]

    target_label_col = _find_column_by_candidates(
        gathered_df, ["Target Label", "Target", "Target Type"]
    )
    normalized_targets: set[str] = set()
    for label in target_labels or []:
        normalized = _normalize_text_value(label)
        if not normalized:
            continue
        normalized_targets.add(normalized)
        if normalized == "cross":
            normalized_targets.add("competitor")
        if normalized == "competitor":
            normalized_targets.add("cross")
    if target_label_col and normalized_targets:
        target_series = data_df[target_label_col].map(_normalize_text_value)
        filtered_df = filtered_df[target_series.isin(normalized_targets)]

    labels = (
        filtered_df[label_col]
        .dropna()
        .astype(str)
        .map(str.strip)
    )
    unique_labels = []
    seen = set()
    for label in labels:
        if not label or label in seen:
            continue
        seen.add(label)
        unique_labels.append(label)
    return unique_labels

def _normalize_text_value(value) -> str:
    if pd.isna(value):
        return ""
    return str(value).strip().lower()


def _normalize_lookup_value(value: str) -> str:
    normalized = str(value).strip().casefold()
    normalized = re.sub(r"[^\w\s]", " ", normalized)
    normalized = normalized.replace("_", " ")
    normalized = " ".join(normalized.split())
    return normalized


def _format_fuzzy_candidates(candidates: list[tuple[float, str]]) -> str:
    return ", ".join(f"{label!r} ({score:.1f})" for score, label in candidates)


def _resolve_column_from_candidates(
    df: pd.DataFrame,
    header_row: pd.Series | None,
    candidates: list[str],
    *,
    context: str,
    threshold: float = 85.0,
) -> tuple[str | None, int, float]:
    column_options: list[tuple[str, object, int]] = []
    for column in df.columns:
        column_options.append((str(column), column, 0))
    if header_row is not None:
        for column, value in header_row.items():
            if pd.isna(value):
                continue
            column_options.append((str(value), column, 1))

    normalized_candidates = [
        (candidate, _normalize_lookup_value(candidate)) for candidate in candidates
    ]
    normalized_options = [
        (label, _normalize_lookup_value(label), column, data_start_idx)
        for label, column, data_start_idx in column_options
        if _normalize_lookup_value(label)
    ]

    exact_matches = []
    for candidate, normalized_candidate in normalized_candidates:
        for label, normalized_label, column, data_start_idx in normalized_options:
            if normalized_candidate == normalized_label:
                exact_matches.append((candidate, label, column, data_start_idx, 100.0))
    if exact_matches:
        unique_columns = {match[2] for match in exact_matches}
        if len(unique_columns) > 1:
            top_candidates = [
                (match[4], match[1]) for match in exact_matches[:5]
            ]
            raise ValueError(
                f"Ambiguous {context} match. Top candidates: "
                f"{_format_fuzzy_candidates(top_candidates)}"
            )
        candidate, label, column, data_start_idx, score = exact_matches[0]
        logger.info(
            'Resolved header "%s" -> "%s" (score %.1f)',
            candidate,
            label,
            score,
        )
        return column, data_start_idx, score

    from difflib import SequenceMatcher

    scored: list[tuple[float, str, object, int, str]] = []
    for candidate, normalized_candidate in normalized_candidates:
        for label, normalized_label, column, data_start_idx in normalized_options:
            score = SequenceMatcher(None, normalized_candidate, normalized_label).ratio() * 100
            scored.append((score, label, column, data_start_idx, candidate))
    if not scored:
        return None, 0, 0.0
    scored.sort(key=lambda item: item[0], reverse=True)
    top_score = scored[0][0]
    if top_score < threshold:
        return None, 0, top_score
    close_matches = [
        (score, label, column)
        for score, label, column, _, _ in scored
        if score >= top_score - 1.0
    ]
    unique_columns = {column for _, _, column in close_matches}
    if len(unique_columns) > 1:
        top_candidates = [(score, label) for score, label, _ in close_matches[:5]]
        raise ValueError(
            f"Ambiguous {context} match. Top candidates: "
            f"{_format_fuzzy_candidates(top_candidates)}"
        )
    score, label, column, data_start_idx, candidate = scored[0]
    logger.info(
        'Resolved header "%s" -> "%s" (score %.1f)',
        candidate,
        label,
        score,
    )
    return column, data_start_idx, score


def _resolve_label_from_text(text: str, labels: list[str], threshold: float = 85.0) -> str:
    normalized_text = _normalize_lookup_value(text)
    if not normalized_text:
        raise ValueError("Slide text is empty after normalization.")
    normalized_labels = []
    exact_matches = []
    for label in labels:
        normalized_label = _normalize_lookup_value(label)
        normalized_labels.append((label, normalized_label))
        if normalized_label == normalized_text:
            exact_matches.append(label)
    if exact_matches:
        resolved = exact_matches[0]
        logger.info(
            "Resolved slide text %r -> %r (score 100.0)",
            text,
            resolved,
        )
        return resolved
    from difflib import SequenceMatcher

    scored = []
    for label, normalized_label in normalized_labels:
        if not normalized_label:
            continue
        score = SequenceMatcher(None, normalized_text, normalized_label).ratio() * 100
        scored.append((score, label))
    scored.sort(reverse=True)
    top_candidates = scored[:5]
    if not scored or scored[0][0] < threshold:
        raise ValueError(
            "No slide text match found (threshold {:.0f}). Top candidates: {}".format(
                threshold,
                _format_fuzzy_candidates(top_candidates),
            )
        )
    best_score, best_label = scored[0]
    if len(scored) > 1 and scored[1][0] >= threshold and (best_score - scored[1][0]) < 2:
        raise ValueError(
            "Ambiguous slide text match (threshold {:.0f}). Top candidates: {}".format(
                threshold,
                _format_fuzzy_candidates(top_candidates),
            )
        )
    logger.info(
        "Resolved slide text %r -> %r (score %.1f)",
        text,
        best_label,
        best_score,
    )
    return best_label


def _resolve_target_level_label_for_slide(slide, labels: list[str]) -> str | None:
    candidates = []
    slide_title = _slide_title(slide)
    if slide_title:
        candidates.append(("title", slide_title))
    slide_name = getattr(slide, "name", None) or ""
    if slide_name:
        candidates.append(("name", slide_name))
    if not candidates:
        return None
    errors = []
    for source, text in candidates:
        try:
            return _resolve_label_from_text(text, labels)
        except ValueError as exc:
            message = str(exc)
            error_message = (
                f"Could not resolve Target Level Label from slide {source} {text!r}: {exc}"
            )
            if message.startswith("No slide text match found") or message.startswith(
                "Slide text is empty after normalization"
            ):
                logger.info("%s", error_message)
                continue
            errors.append(error_message)
    if errors:
        raise ValueError(" | ".join(errors))
    return None


def _normalize_lookup_value(value: str) -> str:
    normalized = str(value).strip().casefold()
    normalized = re.sub(r"[^\w\s]", " ", normalized)
    normalized = normalized.replace("_", " ")
    normalized = " ".join(normalized.split())
    return normalized


def _format_fuzzy_candidates(candidates: list[tuple[float, str]]) -> str:
    return ", ".join(f"{label!r} ({score:.1f})" for score, label in candidates)


def _resolve_label_from_text(text: str, labels: list[str], threshold: float = 85.0) -> str:
    normalized_text = _normalize_lookup_value(text)
    if not normalized_text:
        raise ValueError("Slide text is empty after normalization.")
    normalized_labels = []
    exact_matches = []
    for label in labels:
        normalized_label = _normalize_lookup_value(label)
        normalized_labels.append((label, normalized_label))
        if normalized_label == normalized_text:
            exact_matches.append(label)
    if exact_matches:
        resolved = exact_matches[0]
        logger.info(
            "Resolved slide text %r -> %r (score 100.0)",
            text,
            resolved,
        )
        return resolved
    from difflib import SequenceMatcher

    scored = []
    for label, normalized_label in normalized_labels:
        if not normalized_label:
            continue
        score = SequenceMatcher(None, normalized_text, normalized_label).ratio() * 100
        scored.append((score, label))
    scored.sort(reverse=True)
    top_candidates = scored[:5]
    if not scored or scored[0][0] < threshold:
        raise ValueError(
            "No slide text match found (threshold {:.0f}). Top candidates: {}".format(
                threshold,
                _format_fuzzy_candidates(top_candidates),
            )
        )
    best_score, best_label = scored[0]
    if len(scored) > 1 and scored[1][0] >= threshold and (best_score - scored[1][0]) < 2:
        raise ValueError(
            "Ambiguous slide text match (threshold {:.0f}). Top candidates: {}".format(
                threshold,
                _format_fuzzy_candidates(top_candidates),
            )
        )
    logger.info(
        "Resolved slide text %r -> %r (score %.1f)",
        text,
        best_label,
        best_score,
    )
    return best_label


def _normalize_target_label_value(value: str | None) -> str:
    normalized = _normalize_text_value(value)
    if normalized == "competitor":
        return "cross"
    return normalized


def _target_label_values_from_gathered_df(gathered_df: pd.DataFrame) -> list[str]:
    if gathered_df is None or gathered_df.empty:
        return []
    header_row = gathered_df.iloc[0] if len(gathered_df) else None
    column = _find_column_by_candidates(
        gathered_df, ["Target Label", "Target", "Target Type"]
    )
    data_start_idx = 0
    if not column and header_row is not None:
        column = _find_column_by_row_values(header_row, ["Target Label", "Target", "Target Type"])
        if column:
            data_start_idx = 1
    if not column:
        return []
    labels = (
        gathered_df.iloc[data_start_idx:][column]
        .dropna()
        .astype(str)
        .map(str.strip)
    )
    unique_labels = []
    seen = set()
    for label in labels:
        normalized = _normalize_target_label_value(label)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        unique_labels.append(label)
    return unique_labels


def _resolve_target_label_for_slide(slide, labels: list[str]) -> str | None:
    candidates = []
    slide_title = _slide_title(slide)
    if slide_title:
        candidates.append(("title", slide_title))
    slide_name = getattr(slide, "name", None) or ""
    if slide_name:
        candidates.append(("name", slide_name))
    if not candidates:
        return None
    errors = []
    for source, text in candidates:
        try:
            return _resolve_label_from_text(text, labels)
        except ValueError as exc:
            error_message = f"Could not resolve Target Label from slide {source} {text!r}: {exc}"
            errors.append(error_message)
    if errors:
        raise ValueError(" | ".join(errors))
    return None


def _two_row_column_match(
    group_value: str,
    sub_value: str,
    candidates: list[str],
) -> bool:
    group_key = _normalize_column_name(group_value)
    sub_key = _normalize_column_name(sub_value)
    candidate_keys = {_normalize_column_name(candidate) for candidate in candidates}
    return group_key in candidate_keys or sub_key in candidate_keys


def _parse_two_row_header_dataframe(
    raw_df: pd.DataFrame,
) -> tuple[pd.DataFrame, dict]:
    """Parse a gatheredCN10 file that uses two header rows.

    Returns the data rows with stable internal column IDs plus metadata for UI mapping.

    Example:
        >>> raw = pd.DataFrame(
        ...     [
        ...         ["Promo", "Promo", "", ""],
        ...         ["Feature", "Display", "Target Label", "Year"],
        ...         [1, 2, "Own", 2023],
        ...     ]
        ... )
        >>> data_df, meta = _parse_two_row_header_dataframe(raw)
        >>> meta["group_order"]
        ['Promo']
    """
    if raw_df is None or raw_df.empty or raw_df.shape[0] < 3:
        raise ValueError("The gatheredCN10 file must include two header rows and data rows.")
    header_row1 = raw_df.iloc[0].fillna("")
    header_row2 = raw_df.iloc[1].fillna("")
    columns_meta = []
    group_map: dict[str, list[dict]] = {}
    group_order: list[str] = []
    for idx in range(raw_df.shape[1]):
        group = str(header_row1.iloc[idx]).strip()
        subheader = str(header_row2.iloc[idx]).strip()
        col_id = f"col_{idx}"
        columns_meta.append(
            {
                "id": col_id,
                "group": group,
                "subheader": subheader,
                "position": idx,
            }
        )
        if not group:
            continue
        group_key = _normalize_column_name(group)
        if group_key in {"targetlabel", "year"}:
            continue
        if group not in group_map:
            group_map[group] = []
            group_order.append(group)
        group_map[group].append(
            {
                "id": col_id,
                "subheader": subheader,
                "position": idx,
            }
        )

    target_label_id = None
    year_id = None
    for column in columns_meta:
        if target_label_id is None and _two_row_column_match(
            column["group"],
            column["subheader"],
            ["Target Label"],
        ):
            target_label_id = column["id"]
        if year_id is None and _two_row_column_match(
            column["group"],
            column["subheader"],
            ["Year"],
        ):
            year_id = column["id"]

    data_df = raw_df.iloc[2:].reset_index(drop=True).copy()
    data_df.columns = [col["id"] for col in columns_meta]
    metadata = {
        "columns": columns_meta,
        "groups": group_map,
        "group_order": group_order,
        "target_label_id": target_label_id,
        "year_id": year_id,
    }
    return data_df, metadata


def _unique_column_values(data_df: pd.DataFrame, column_id: str) -> list[str]:
    if column_id not in data_df.columns:
        return []
    values = (
        data_df[column_id]
        .dropna()
        .astype(str)
        .map(str.strip)
    )
    unique_values = []
    seen = set()
    for value in values:
        if not value or value in seen:
            continue
        seen.add(value)
        unique_values.append(value)
    return unique_values


def _compute_bucket_deltas(
    data_df: pd.DataFrame,
    metadata: dict,
    bucket_config: dict[str, dict[str, list[str]]],
    year1: str,
    year2: str,
) -> list[tuple[str, float]]:
    """Compute Year2-Year1 deltas for each bucket group.

    bucket_config maps group -> {"target_labels": [...], "subheaders_included": [...]}
    """
    target_label_id = metadata.get("target_label_id")
    year_id = metadata.get("year_id")
    if not target_label_id:
        raise ValueError("The gatheredCN10 file is missing the Target Label column.")
    if not year_id:
        raise ValueError("The gatheredCN10 file is missing the Year column.")

    normalized_year1 = _normalize_text_value(year1)
    normalized_year2 = _normalize_text_value(year2)

    target_series = data_df[target_label_id].map(_normalize_text_value)
    year_series = data_df[year_id].map(_normalize_text_value)

    deltas: list[tuple[str, float]] = []
    group_order = metadata.get("group_order", [])
    ordered_groups = [group for group in group_order if group in bucket_config]
    if not ordered_groups:
        ordered_groups = list(bucket_config.keys())
    for group in ordered_groups:
        config = bucket_config.get(group, {})
        selected_cols = [
            col for col in config.get("subheaders_included", []) if col in data_df.columns
        ]
        target_labels = config.get("target_labels")
        if target_labels is None:
            target_labels = []
        if not target_labels:
            continue
        ordered_targets = []
        normalized_targets = []
        for label in target_labels:
            normalized = _normalize_text_value(label)
            if normalized and normalized not in normalized_targets:
                normalized_targets.append(normalized)
                ordered_targets.append((label, normalized))
        target_label_sequence = []
        if "own" in normalized_targets:
            target_label_sequence.append(("Own", "own"))
        if "cross" in normalized_targets:
            target_label_sequence.append(("Cross", "cross"))
        for label, normalized in ordered_targets:
            if normalized not in {"own", "cross"}:
                target_label_sequence.append((label, normalized))
        if not target_label_sequence:
            deltas.append((group, 0.0))
            continue
        if not selected_cols:
            for label, _ in target_label_sequence:
                display_label = DISPLAY_LABEL.get(label, label)
                deltas.append((f"{display_label} {group}", 0.0))
            continue
        values_df = data_df[selected_cols].apply(pd.to_numeric, errors="coerce").fillna(0)
        year1_mask = year_series == normalized_year1
        year2_mask = year_series == normalized_year2
        for label, normalized in target_label_sequence:
            target_mask = target_series == normalized
            year1_sum = values_df[year1_mask & target_mask].sum().sum()
            year2_sum = values_df[year2_mask & target_mask].sum().sum()
            display_label = DISPLAY_LABEL.get(label, label)
            deltas.append((f"{display_label} {group}", float(year2_sum - year1_sum)))
    return deltas

def _resolve_base_value_columns(gathered_df: pd.DataFrame) -> tuple[dict, int]:
    column_candidates = {
        "target_level": ["Target Level Label", "Target Level"],
        "target_label": ["Target Label", "Target", "Target Type"],
        "year": ["Year", "Model Year"],
        "actuals": ["Actuals", "Actual"],
    }
    columns = {}
    data_start_idx = 0
    header_row = gathered_df.iloc[0] if len(gathered_df) else None
    for key, candidates in column_candidates.items():
        column = _find_column_by_candidates(gathered_df, candidates)
        if not column and header_row is not None:
            column = _find_column_by_row_values(header_row, candidates)
            if column:
                data_start_idx = 1
        if not column:
            raise ValueError(
                "The gatheredCN10 file is missing the "
                f"{' / '.join(candidates)} column needed for the waterfall base."
            )
        columns[key] = column
    if columns["target_level"] == columns["target_label"]:
        alt_candidates = ["Target Level Label", "Target Level"]
        column = _find_column_by_candidates(gathered_df, alt_candidates)
        if not column and header_row is not None:
            column = _find_column_by_row_values(header_row, alt_candidates)
            if column:
                data_start_idx = 1
        if column and column != columns["target_label"]:
            columns["target_level"] = column
        else:
            raise ValueError(
                "The gatheredCN10 file needs separate Target Level Label and Target Label columns."
            )
    return columns, data_start_idx


def _waterfall_base_values(
    gathered_df: pd.DataFrame,
    target_level_label: str,
    year1: str | None = None,
    year2: str | None = None,
) -> tuple[float, float]:
    if gathered_df is None or gathered_df.empty:
        raise ValueError("The gatheredCN10 file is empty.")
    columns, data_start_idx = _resolve_base_value_columns(gathered_df)
    data_df = gathered_df.iloc[data_start_idx:]
    target_level = _normalize_text_value(target_level_label)
    target_level_series = data_df[columns["target_level"]].map(_normalize_text_value)
    target_label_series = data_df[columns["target_label"]].map(_normalize_text_value)
    year_series = data_df[columns["year"]].map(_normalize_text_value)
    actuals = pd.to_numeric(data_df[columns["actuals"]], errors="coerce").fillna(0)
    base_filter = (target_level_series == target_level) & (target_label_series == "own")
    normalized_year1 = _normalize_text_value(year1) if year1 is not None else "year1"
    normalized_year2 = _normalize_text_value(year2) if year2 is not None else "year2"
    year1_total = actuals[base_filter & (year_series == normalized_year1)].sum()
    year2_total = actuals[base_filter & (year_series == normalized_year2)].sum()
    return year1_total, year2_total


def _format_lab_base_value(value: float) -> str:
    if value is None or pd.isna(value):
        return ""
    abs_value = abs(value)
    if abs_value >= 1_000_000:
        scaled = value / 1_000_000
        suffix = "m"
    elif abs_value >= 1_000:
        scaled = value / 1_000
        suffix = "k"
    else:
        return str(int(value)) if float(value).is_integer() else str(value)
    formatted = f"{scaled:g}"
    return f"{formatted}{suffix}"


def _load_chart_workbook(chart):
    xlsx_blob = chart.part.chart_workbook.xlsx_part.blob
    return load_workbook(io.BytesIO(xlsx_blob))


def _save_chart_workbook(chart, workbook) -> None:
    stream = io.BytesIO()
    workbook.save(stream)
    chart.part.chart_workbook.xlsx_part.blob = stream.getvalue()


def _chart_namespace_map(root) -> dict:
    nsmap = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
             "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
             "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
             "c15": "http://schemas.microsoft.com/office/drawing/2012/chart"}
    for prefix, uri in (root.nsmap or {}).items():
        if prefix and uri:
            nsmap[prefix] = uri
    return nsmap


def _range_values_from_worksheet(ws, ref: str) -> list[list]:
    if not ref:
        return []
    normalized = str(ref)
    if "!" in normalized:
        _, normalized = normalized.split("!", 1)
    normalized = normalized.replace("$", "")
    try:
        min_col, min_row, max_col, max_row = range_boundaries(normalized)
    except ValueError:
        return []
    rows = []
    for row_idx in range(min_row, max_row + 1):
        row = []
        for col_idx in range(min_col, max_col + 1):
            row.append(ws.cell(row=row_idx, column=col_idx).value)
        rows.append(row)
    return rows


def _range_cells_from_worksheet(ws, ref: str) -> list[list]:
    if not ref:
        return []
    normalized = str(ref)
    if "!" in normalized:
        _, normalized = normalized.split("!", 1)
    normalized = normalized.replace("$", "")
    try:
        min_col, min_row, max_col, max_row = range_boundaries(normalized)
    except ValueError:
        return []
    rows = []
    for row_idx in range(min_row, max_row + 1):
        row = []
        for col_idx in range(min_col, max_col + 1):
            row.append(ws.cell(row=row_idx, column=col_idx))
        rows.append(row)
    return rows


def _worksheet_and_range_from_formula(workbook, formula: str) -> tuple:
    if not formula:
        return workbook.active, "", None
    sheet_name = None
    ref = str(formula)
    if "!" in ref:
        match = re.match(
            r"^(?P<sheet>(?:'[^']*(?:''[^']*)*'|[^!]+))!(?P<ref>.+)$",
            ref,
        )
        if match:
            sheet_part = match.group("sheet")
            ref = match.group("ref")
            if sheet_part.startswith("'") and sheet_part.endswith("'"):
                sheet_name = sheet_part[1:-1].replace("''", "'")
            else:
                sheet_name = sheet_part
        else:
            sheet_part, ref = ref.split("!", 1)
            sheet_name = sheet_part.strip("'")
    ref = ref.replace("$", "")
    ws = workbook.active
    if sheet_name:
        if sheet_name not in workbook.sheetnames:
            resolved = _find_sheet_by_candidates(workbook.sheetnames, sheet_name)
            if resolved:
                logger.info(
                    "Chart cache: resolved sheet '%s' -> '%s' from formula '%s'",
                    sheet_name,
                    resolved,
                    formula,
                )
                sheet_name = resolved
            else:
                raise ValueError(
                    f"Chart cache: sheet '{sheet_name}' from formula '{formula}' not found."
                )
        ws = workbook[sheet_name]
    return ws, ref, sheet_name


def _format_sheet_reference(sheet_name: str) -> str:
    if sheet_name is None:
        return ""
    if sheet_name.startswith("'") and sheet_name.endswith("'"):
        return sheet_name
    if re.search(r"[^A-Za-z0-9_]", sheet_name):
        return f"'{sheet_name}'"
    return sheet_name


def _build_cell_range_formula(sheet_name: str | None, col_idx: int, row_start: int, row_end: int) -> str:
    col_letter = get_column_letter(col_idx)
    sheet_prefix = ""
    if sheet_name:
        sheet_prefix = f"{_format_sheet_reference(sheet_name)}!"
    return f"{sheet_prefix}${col_letter}${row_start}:${col_letter}${row_end}"


def _range_boundaries_from_formula(formula: str) -> tuple[int, int, int, int] | None:
    if not formula:
        return None
    ref = str(formula)
    if "!" in ref:
        _, ref = ref.split("!", 1)
    ref = ref.replace("$", "")
    try:
        return range_boundaries(ref)
    except ValueError:
        return None


def _flatten_cell_values(values: list[list]) -> list:
    if not values:
        return []
    if len(values) == 1:
        return list(values[0])
    if all(len(row) == 1 for row in values):
        return [row[0] for row in values]
    flattened = []
    for row in values:
        flattened.extend(row)
    return flattened


def _log_num_cache_warning(
    value,
    fallback: float,
    sheet_name: str | None,
    cell_ref: str | None,
) -> None:
    global _num_cache_warning_count
    if _num_cache_warning_count >= _NUM_CACHE_WARNING_LIMIT:
        return
    _num_cache_warning_count += 1
    location = []
    if sheet_name:
        location.append(f"sheet={sheet_name}")
    if cell_ref:
        location.append(f"cell={cell_ref}")
    location_text = f" ({', '.join(location)})" if location else ""
    logger.warning(
        "Chart cache: coerced non-numeric value%s %r to %s",
        location_text,
        value,
        fallback,
    )


def safe_float(
    value,
    *,
    sheet_name: str | None = None,
    cell_ref: str | None = None,
) -> float:
    if value is None:
        _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
        return 0.0
    if isinstance(value, str):
        stripped = value.strip()
        if stripped == "":
            _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
            return 0.0
        if stripped.startswith("="):
            _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
            return 0.0
        normalized = stripped.replace(",", "")
        if re.search(r"[A-Za-z]", normalized):
            _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
            return 0.0
        if not re.fullmatch(r"[+-]?(?:\d+(\.\d*)?|\.\d+)", normalized):
            _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
            return 0.0
        try:
            return float(normalized)
        except (TypeError, ValueError):
            _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
            return 0.0
    if isinstance(value, numbers.Real):
        if pd.isna(value):
            _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
            return 0.0
        return float(value)
    try:
        converted = float(value)
    except (TypeError, ValueError):
        _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
        return 0.0
    if pd.isna(converted):
        _log_num_cache_warning(value, 0.0, sheet_name, cell_ref)
        return 0.0
    return converted


def _all_blank(values: list) -> bool:
    if not values:
        return True
    for value in values:
        if value is None:
            continue
        if isinstance(value, str) and value.strip() == "":
            continue
        return False
    return True


def _ensure_str_cache(str_ref) -> tuple:
    str_cache = str_ref.find("c:strCache", namespaces=_chart_namespace_map(str_ref))
    created = False
    if str_cache is None:
        str_cache = etree.SubElement(
            str_ref,
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}strCache",
        )
        created = True
    return str_cache, created


def _update_num_cache(num_cache, values: list) -> None:
    if num_cache is None:
        return
    pt_count = num_cache.find("c:ptCount", namespaces=_chart_namespace_map(num_cache))
    if pt_count is None:
        pt_count = etree.SubElement(
            num_cache,
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}ptCount",
        )
    pt_count.set("val", str(len(values)))
    for pt in list(num_cache.findall("c:pt", namespaces=_chart_namespace_map(num_cache))):
        num_cache.remove(pt)
    for idx, value in enumerate(values):
        cell = value if hasattr(value, "value") and hasattr(value, "coordinate") else None
        raw_value = value.value if cell is not None else value
        sheet_name = cell.parent.title if cell is not None else None
        cell_ref = cell.coordinate if cell is not None else None
        normalized_value = safe_float(
            raw_value,
            sheet_name=sheet_name,
            cell_ref=cell_ref,
        )
        pt = etree.SubElement(
            num_cache,
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}pt",
            idx=str(idx),
        )
        v = etree.SubElement(
            pt, "{http://schemas.openxmlformats.org/drawingml/2006/chart}v"
        )
        v.text = str(normalized_value)


def _update_str_cache(str_cache, values: list[str]) -> None:
    if str_cache is None:
        return
    pt_count = str_cache.find("c:ptCount", namespaces=_chart_namespace_map(str_cache))
    if pt_count is None:
        pt_count = etree.SubElement(
            str_cache,
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}ptCount",
        )
    pt_count.set("val", str(len(values)))
    for pt in list(str_cache.findall("c:pt", namespaces=_chart_namespace_map(str_cache))):
        str_cache.remove(pt)
    for idx, value in enumerate(values):
        pt = etree.SubElement(
            str_cache,
            "{http://schemas.openxmlformats.org/drawingml/2006/chart}pt",
            idx=str(idx),
        )
        v = etree.SubElement(
            pt, "{http://schemas.openxmlformats.org/drawingml/2006/chart}v"
        )
        v.text = "" if value is None else str(value)


def _update_c15_label_range_cache(
    container,
    formula: str | None,
    labels: list[str],
    nsmap: dict,
    label_context: str,
) -> int:
    c15_blocks = []
    if container is not None and str(getattr(container, "tag", "")).endswith("datalabelsRange"):
        c15_blocks.append(container)
    c15_blocks += container.findall(".//c15:datalabelsRange", namespaces=nsmap)
    if not c15_blocks:
        try:
            c15_blocks = container.xpath(".//*[local-name()='datalabelsRange']")
        except Exception:
            c15_blocks = []
    if c15_blocks:
        seen = set()
        deduped = []
        for block in c15_blocks:
            if id(block) in seen:
                continue
            seen.add(id(block))
            deduped.append(block)
        c15_blocks = deduped
    if not c15_blocks:
        logger.info(
            "Waterfall chart cache update: %s no c15 label-range block found",
            label_context,
        )
        return 0
    logger.info(
        "Waterfall chart cache update: %s c15 label-range blocks found %s",
        label_context,
        len(c15_blocks),
    )
    for c15_block in c15_blocks:
        block_ns = etree.QName(c15_block).namespace or nsmap.get(
            "c15", "http://schemas.microsoft.com/office/drawing/2012/chart"
        )
        f_node = c15_block.find("c15:f", namespaces=nsmap)
        if f_node is None:
            try:
                f_node = c15_block.xpath("./*[local-name()='f']")[0]
            except Exception:
                f_node = None
        if f_node is None:
            f_node = etree.SubElement(c15_block, f"{{{block_ns}}}f")
        if formula:
            f_node.text = formula
            logger.info(
                "Waterfall chart cache update: %s c15 label-range formula set to %s",
                label_context,
                formula,
            )
        cache = c15_block.find("c15:dlblRangeCache", namespaces=nsmap)
        if cache is None:
            try:
                cache = c15_block.xpath("./*[local-name()='dlblRangeCache']")[0]
            except Exception:
                cache = None
        if cache is None:
            cache = etree.SubElement(c15_block, f"{{{block_ns}}}dlblRangeCache")
        pt_count = cache.find("c15:ptCount", namespaces=nsmap)
        if pt_count is None:
            try:
                pt_count = cache.xpath("./*[local-name()='ptCount']")[0]
            except Exception:
                pt_count = None
        if pt_count is None:
            pt_count = etree.SubElement(cache, f"{{{block_ns}}}ptCount")
        pt_count.set("val", str(len(labels)))
        for pt in list(cache.findall("c15:pt", namespaces=nsmap)):
            cache.remove(pt)
        for pt in list(cache.xpath("./*[local-name()='pt']")):
            cache.remove(pt)
        for idx, value in enumerate(labels):
            pt = etree.SubElement(cache, f"{{{block_ns}}}pt", idx=str(idx))
            v = etree.SubElement(pt, f"{{{block_ns}}}v")
            v.text = "" if value is None else str(value)
        logger.info(
            "Waterfall chart cache update: %s c15 label-range cached %s points",
            label_context,
            len(labels),
        )
    return len(c15_blocks)


def _update_waterfall_chart_caches(chart, workbook, categories: list[str]) -> None:
    chart_part = chart.part
    root = chart_part._element
    nsmap = _chart_namespace_map(root)
    ws = workbook.active
    label_columns = {
        col_idx: ws.cell(row=1, column=col_idx).value
        for col_idx in range(1, ws.max_column + 1)
        if ws.cell(row=1, column=col_idx).value
        and _normalize_column_name(str(ws.cell(row=1, column=col_idx).value)).startswith("labs")
    }
    if label_columns:
        logger.info(
            "Waterfall chart cache update: label columns found %s",
            {idx: str(value) for idx, value in label_columns.items()},
        )
    categories_values = ["" if value is None else str(value) for value in categories]
    categories_count = len(categories_values)
    logger.info("Waterfall chart cache update: %s category points", categories_count)

    series_names = [series.name for series in chart.series]
    series_point_counts: dict[int, int] = {}
    series_category_bounds: dict[int, tuple[int, int, str | None]] = {}
    series_value_bounds: dict[int, tuple[int, int, str | None]] = {}

    for idx, ser in enumerate(root.findall(".//c:ser", namespaces=nsmap), start=1):
        num_ref = ser.find("c:val/c:numRef", namespaces=nsmap)
        if num_ref is None:
            continue
        f_node = num_ref.find("c:f", namespaces=nsmap)
        if f_node is None or not f_node.text:
            continue
        value_ws, value_ref, _ = _worksheet_and_range_from_formula(workbook, f_node.text)
        value_rows = _range_cells_from_worksheet(value_ws, value_ref)
        series_values = _flatten_cell_values(value_rows)
        num_cache = num_ref.find("c:numCache", namespaces=nsmap)
        _update_num_cache(num_cache, series_values)
        series_point_counts[idx] = len(series_values)
        bounds = _range_boundaries_from_formula(f_node.text)
        if bounds:
            _, min_row, _, max_row = bounds
            series_value_bounds[idx] = (min_row, max_row, value_ws.title)
        logger.info(
            "Waterfall chart cache update: series %s cached %s points",
            idx,
            len(series_values),
        )

    category_cache_updates = 0
    category_cache_count = None
    for idx, ser in enumerate(root.findall(".//c:ser", namespaces=nsmap), start=1):
        series_label = series_names[idx - 1] if idx - 1 < len(series_names) else f"Series {idx}"
        cat_node = ser.find("c:cat", namespaces=nsmap)
        if cat_node is None:
            logger.info(
                "Waterfall chart cache update: series %s category ref not found",
                series_label,
            )
            continue
        cat_ref = cat_node.find("c:strRef", namespaces=nsmap)
        cat_ref_type = "strRef"
        num_ref = None
        if cat_ref is None:
            num_ref = cat_node.find("c:numRef", namespaces=nsmap)
            cat_ref_type = "numRef"
            cat_ref = num_ref
        if cat_ref is None:
            logger.info(
                "Waterfall chart cache update: series %s category ref not found",
                series_label,
            )
            continue
        f_node = cat_ref.find("c:f", namespaces=nsmap)
        if f_node is None or not f_node.text:
            logger.info(
                "Waterfall chart cache update: series %s category ref formula missing",
                series_label,
            )
            continue
        logger.info(
            "Waterfall chart cache update: series %s category ref type %s formula %s",
            series_label,
            cat_ref_type,
            f_node.text,
        )
        if cat_ref_type == "numRef" and num_ref is not None:
            f_text = f_node.text
            num_ref_index = list(cat_node).index(num_ref)
            cat_node.remove(num_ref)
            cat_ref = etree.Element("{http://schemas.openxmlformats.org/drawingml/2006/chart}strRef")
            f_node = etree.SubElement(
                cat_ref, "{http://schemas.openxmlformats.org/drawingml/2006/chart}f"
            )
            f_node.text = f_text
            cat_node.insert(num_ref_index, cat_ref)
            cat_ref_type = "strRef"
        logger.info(
            "Waterfall chart cache update: series %s category ref formula %s",
            series_label,
            f_node.text,
        )
        cat_ws, cat_ref_range, cat_sheet = _worksheet_and_range_from_formula(
            workbook, f_node.text
        )
        category_rows = _range_values_from_worksheet(cat_ws, cat_ref_range)
        category_values = _flatten_cell_values(category_rows)
        if _all_blank(category_values):
            raise ValueError(
                f"Chart cache: category range '{f_node.text}' for series '{series_label}' is blank."
            )
        if not category_values and categories_values:
            category_values = categories_values
        category_values = ["" if value is None else str(value) for value in category_values]
        bounds = _range_boundaries_from_formula(f_node.text)
        if bounds:
            _, min_row, _, max_row = bounds
            series_category_bounds[idx] = (min_row, max_row, cat_sheet or cat_ws.title)
        str_cache, created = _ensure_str_cache(cat_ref)
        logger.info(
            "Waterfall chart cache update: series %s category strCache %s",
            series_label,
            "created" if created else "existing",
        )
        _update_str_cache(str_cache, category_values)
        category_cache_updates += 1
        category_cache_count = len(category_values)
        logger.info(
            "Waterfall chart cache update: series %s cached %s category points",
            series_label,
            len(category_values),
        )
    if category_cache_updates:
        logger.info(
            "Waterfall chart cache update: %s category cache points",
            category_cache_count if category_cache_count is not None else categories_count,
        )

    label_cache_updates = 0
    label_cache_missing = 0
    c15_label_updates = 0
    def _collect_label_refs():
        series_refs = []
        for idx, ser in enumerate(root.findall(".//c:ser", namespaces=nsmap), start=1):
            series_label = series_names[idx - 1] if idx - 1 < len(series_names) else f"Series {idx}"
            chart_series = chart.series[idx - 1] if idx - 1 < len(chart.series) else None
            label_refs = ser.findall(".//c:dLbls//c:dLbl//c:tx//c:strRef", namespaces=nsmap)
            label_refs += ser.findall(".//c:dLbls//c:tx//c:strRef", namespaces=nsmap)
            series_refs.append((idx, series_label, chart_series, label_refs, ser))
        plot_level_refs = root.findall(
            "c:plotArea//c:dLbls//c:tx//c:strRef", namespaces=nsmap
        )
        plot_level_dlbls = root.findall("c:plotArea/c:dLbls", namespaces=nsmap)
        plot_refs = []
        for ref_node in plot_level_refs:
            current = ref_node.getparent()
            has_series_ancestor = False
            while current is not None:
                if current.tag.endswith("ser"):
                    has_series_ancestor = True
                    break
                current = current.getparent()
            if not has_series_ancestor:
                plot_refs.append(ref_node)
        series_ref_count_local = sum(len(entry[3]) for entry in series_refs)
        return series_refs, plot_refs, series_ref_count_local, plot_level_dlbls

    series_label_refs, plot_only_refs, series_ref_count, plot_level_dlbls = _collect_label_refs()
    logger.info(
        "Waterfall chart cache update: data label strRef nodes found (series=%s, plot-level=%s)",
        series_ref_count,
        len(plot_only_refs),
    )
    if series_ref_count == 0 and (plot_only_refs or plot_level_dlbls):
        plot_dlbls = None
        if plot_level_dlbls:
            plot_dlbls = plot_level_dlbls[0]
        for ref_node in plot_only_refs:
            current = ref_node
            while current is not None and not current.tag.endswith("dLbls"):
                current = current.getparent()
            if current is not None:
                plot_dlbls = current
                break
        if plot_dlbls is None:
            plot_dlbls = root.find(".//c:plotArea//c:dLbls", namespaces=nsmap)
        if plot_dlbls is None:
            logger.info(
                "Waterfall chart cache update: plot-level data labels not found for promotion",
            )
        else:
            promoted = 0
            for ser in root.findall(".//c:plotArea//c:ser", namespaces=nsmap):
                if ser.find("c:dLbls", namespaces=nsmap) is None:
                    ser.append(copy.deepcopy(plot_dlbls))
                    promoted += 1
            parent = plot_dlbls.getparent()
            if parent is not None:
                parent.remove(plot_dlbls)
            logger.info(
                "Waterfall chart cache update: promoted plot-level data labels into %s series and removed plot-level node",
                promoted,
            )
            series_label_refs, plot_only_refs, series_ref_count, plot_level_dlbls = _collect_label_refs()
            logger.info(
                "Waterfall chart cache update: data label strRef nodes found after promotion (series=%s, plot-level=%s)",
                series_ref_count,
                len(plot_only_refs),
            )
    for idx, series_label, chart_series, label_refs, ser in series_label_refs:
        seen_refs = set()
        deduped_refs = []
        for ref in label_refs:
            if id(ref) in seen_refs:
                continue
            seen_refs.add(id(ref))
            deduped_refs.append(ref)
        mapped_header, mapped_labs_column = _resolve_waterfall_labs_column(ws, series_label)
        if mapped_labs_column is None:
            value_formula_node = ser.find("c:val/c:numRef/c:f", namespaces=nsmap)
            if value_formula_node is not None and value_formula_node.text:
                value_bounds = _range_boundaries_from_formula(value_formula_node.text)
                if value_bounds:
                    value_col = value_bounds[0]
                    base_header = ws.cell(row=1, column=value_col).value
                    if base_header:
                        derived_header = f"labs-{str(base_header).strip()}"
                        derived_col = _find_header_column(ws, [derived_header])
                        if derived_col is not None:
                            resolved_header = ws.cell(row=1, column=derived_col).value
                            mapped_header = derived_header
                            mapped_labs_column = derived_col
                            logger.info(
                                "Waterfall chart cache update: series %s fallback labs header resolved %s -> %s (%s)",
                                series_label,
                                derived_header,
                                resolved_header,
                                get_column_letter(derived_col),
                            )
        if mapped_header and mapped_labs_column:
            logger.info(
                "Waterfall chart cache update: series %s mapped to %s (%s)",
                series_label,
                mapped_header,
                get_column_letter(mapped_labs_column),
            )
        else:
            logger.info(
                "Waterfall chart cache update: series %s labs mapping missing",
                series_label,
            )
        bounds = series_category_bounds.get(idx) or series_value_bounds.get(idx)
        if bounds:
            min_row, max_row, sheet_name = bounds
        else:
            value_formula_node = ser.find("c:val/c:numRef/c:f", namespaces=nsmap)
            min_row = max_row = None
            sheet_name = None
            if value_formula_node is not None and value_formula_node.text:
                value_bounds = _range_boundaries_from_formula(value_formula_node.text)
                if value_bounds:
                    _, min_row, _, max_row = value_bounds
                    _, _, sheet_name = _worksheet_and_range_from_formula(
                        workbook, value_formula_node.text
                    )
        c15_updated = False
        if deduped_refs:
            for ref_node in deduped_refs:
                f_node = ref_node.find("c:f", namespaces=nsmap)
                if f_node is None or not f_node.text:
                    logger.info(
                        "Waterfall chart cache update: series %s data label ref formula missing",
                        series_label,
                    )
                    continue
                original_formula = f_node.text
                logger.info(
                    "Waterfall chart cache update: series %s data label ref formula %s",
                    series_label,
                    f_node.text,
                )
                expected_formula = None
                formula_bounds = _range_boundaries_from_formula(f_node.text)
                formula_col = formula_bounds[0] if formula_bounds else None
                effective_label_col = mapped_labs_column or formula_col
                if mapped_labs_column and min_row is not None and max_row is not None:
                    expected_formula = _build_cell_range_formula(
                        sheet_name,
                        effective_label_col,
                        min_row,
                        max_row,
                    )
                    if f_node.text != expected_formula:
                        logger.info(
                            "Waterfall chart cache update: series %s label formula %s -> %s",
                            series_label,
                            f_node.text,
                            expected_formula,
                        )
                        f_node.text = expected_formula
                elif effective_label_col and min_row is not None and max_row is not None:
                    expected_formula = _build_cell_range_formula(
                        sheet_name,
                        effective_label_col,
                        min_row,
                        max_row,
                    )
                elif effective_label_col:
                    series_points = series_point_counts.get(idx)
                    if series_points:
                        expected_formula = _build_cell_range_formula(
                            sheet_name or ws.title,
                            effective_label_col,
                            2,
                            1 + series_points,
                        )
                        if f_node.text != expected_formula:
                            logger.info(
                                "Waterfall chart cache update: series %s label formula %s -> %s",
                                series_label,
                                f_node.text,
                                expected_formula,
                            )
                            f_node.text = expected_formula
                logger.info(
                    "Waterfall chart cache update: series %s label formula resolved %s -> %s",
                    series_label,
                    original_formula,
                    f_node.text,
                )
                label_ws, label_ref_range, _ = _worksheet_and_range_from_formula(
                    workbook, f_node.text
                )
                label_rows = _range_values_from_worksheet(label_ws, label_ref_range)
                label_values = _flatten_cell_values(label_rows)
                if _all_blank(label_values):
                    logger.info(
                        "Waterfall chart cache update: series %s data label range '%s' is blank; skipping cache update",
                        series_label,
                        f_node.text,
                    )
                    continue
                series_points = series_point_counts.get(idx, len(label_values))
                if len(label_values) < series_points:
                    label_values += ["" for _ in range(series_points - len(label_values))]
                elif len(label_values) > series_points:
                    label_values = label_values[:series_points]
                str_cache, created = _ensure_str_cache(ref_node)
                if created:
                    label_cache_missing += 1
                _update_str_cache(
                    str_cache,
                    ["" if value is None else str(value) for value in label_values],
                )
                label_cache_updates += 1
                if not c15_updated:
                    d_lbls_node = ref_node
                    while d_lbls_node is not None and not d_lbls_node.tag.endswith("dLbls"):
                        d_lbls_node = d_lbls_node.getparent()
                    c15_container = d_lbls_node if d_lbls_node is not None else ser
                    c15_label_updates += _update_c15_label_range_cache(
                        c15_container,
                        expected_formula or f_node.text,
                        ["" if value is None else str(value) for value in label_values],
                        nsmap,
                        f"series {series_label}",
                    )
                    c15_updated = True
                logger.info(
                    "Waterfall chart cache update: series %s cached %s data label points",
                    series_label,
                    len(label_values),
                )
                if expected_formula:
                    logger.info(
                        "Waterfall chart cache update: series %s data label ref updated to %s",
                        series_label,
                        expected_formula,
                    )
                logger.info(
                    "Waterfall chart cache update: series %s data label formula now %s",
                    series_label,
                    f_node.text,
                )
                logger.info(
                    "Waterfall chart cache update: series %s data label formula %s cached %s points",
                    series_label,
                    f_node.text,
                    len(label_values),
                )
        if not deduped_refs:
            c15_ranges = ser.findall(".//c15:datalabelsRange", namespaces=nsmap)
            if not c15_ranges:
                try:
                    c15_ranges = ser.xpath(".//*[local-name()='datalabelsRange']")
                except Exception:
                    c15_ranges = []
            if c15_ranges:
                logger.info(
                    "Waterfall chart cache update: series %s c15 label ranges found without c:strRef",
                    series_label,
                )
            for c15_range in c15_ranges:
                c15_formula_node = c15_range.find("c15:f", namespaces=nsmap)
                if c15_formula_node is None:
                    try:
                        c15_formula_node = c15_range.xpath("./*[local-name()='f']")[0]
                    except Exception:
                        c15_formula_node = None
                if c15_formula_node is None or not c15_formula_node.text:
                    logger.info(
                        "Waterfall chart cache update: series %s c15 label formula missing",
                        series_label,
                    )
                    continue
                original_c15_formula = c15_formula_node.text
                logger.info(
                    "Waterfall chart cache update: series %s c15 label formula %s",
                    series_label,
                    c15_formula_node.text,
                )
                expected_formula = None
                formula_bounds = _range_boundaries_from_formula(c15_formula_node.text)
                formula_col = formula_bounds[0] if formula_bounds else None
                effective_label_col = mapped_labs_column or formula_col
                if effective_label_col and min_row is not None and max_row is not None:
                    expected_formula = _build_cell_range_formula(
                        sheet_name,
                        effective_label_col,
                        min_row,
                        max_row,
                    )
                elif effective_label_col:
                    series_points = series_point_counts.get(idx)
                    if series_points:
                        expected_formula = _build_cell_range_formula(
                            sheet_name or ws.title,
                            effective_label_col,
                            2,
                            1 + series_points,
                        )
                if expected_formula and c15_formula_node.text != expected_formula:
                    logger.info(
                        "Waterfall chart cache update: series %s c15 label formula %s -> %s",
                        series_label,
                        c15_formula_node.text,
                        expected_formula,
                    )
                    c15_formula_node.text = expected_formula
                logger.info(
                    "Waterfall chart cache update: series %s c15 label formula resolved %s -> %s",
                    series_label,
                    original_c15_formula,
                    c15_formula_node.text,
                )
                label_ws, label_ref_range, _ = _worksheet_and_range_from_formula(
                    workbook, c15_formula_node.text
                )
                label_rows = _range_values_from_worksheet(label_ws, label_ref_range)
                label_values = _flatten_cell_values(label_rows)
                if _all_blank(label_values):
                    logger.info(
                        "Waterfall chart cache update: series %s c15 label range '%s' is blank; skipping cache update",
                        series_label,
                        c15_formula_node.text,
                    )
                    continue
                series_points = series_point_counts.get(idx, len(label_values))
                if len(label_values) < series_points:
                    label_values += ["" for _ in range(series_points - len(label_values))]
                elif len(label_values) > series_points:
                    label_values = label_values[:series_points]
                c15_label_updates += _update_c15_label_range_cache(
                    c15_range,
                    expected_formula or c15_formula_node.text,
                    ["" if value is None else str(value) for value in label_values],
                    nsmap,
                    f"series {series_label}",
                )
                logger.info(
                    "Waterfall chart cache update: series %s c15 label formula now %s",
                    series_label,
                    c15_formula_node.text,
                )
    for ref_node in plot_only_refs:
        f_node = ref_node.find("c:f", namespaces=nsmap)
        if f_node is None or not f_node.text:
            logger.info(
                "Waterfall chart cache update: plot-level data label ref formula missing",
            )
            continue
        logger.info(
            "Waterfall chart cache update: plot-level data label ref formula %s",
            f_node.text,
        )
        label_ws, label_ref_range, _ = _worksheet_and_range_from_formula(
            workbook, f_node.text
        )
        label_rows = _range_values_from_worksheet(label_ws, label_ref_range)
        label_values = _flatten_cell_values(label_rows)
        if _all_blank(label_values):
            raise ValueError(
                f"Chart cache: plot-level data label range '{f_node.text}' is blank."
            )
        series_points = categories_count or len(label_values)
        if len(label_values) < series_points:
            label_values += ["" for _ in range(series_points - len(label_values))]
        elif len(label_values) > series_points:
            label_values = label_values[:series_points]
        str_cache, created = _ensure_str_cache(ref_node)
        if created:
            label_cache_missing += 1
        _update_str_cache(
            str_cache,
            ["" if value is None else str(value) for value in label_values],
        )
        label_cache_updates += 1
        d_lbls_node = ref_node
        while d_lbls_node is not None and not d_lbls_node.tag.endswith("dLbls"):
            d_lbls_node = d_lbls_node.getparent()
        if d_lbls_node is not None:
            c15_label_updates += _update_c15_label_range_cache(
                d_lbls_node,
                f_node.text,
                ["" if value is None else str(value) for value in label_values],
                nsmap,
                "plot-level",
            )
        logger.info(
            "Waterfall chart cache update: plot-level cached %s data label points",
            len(label_values),
        )
        logger.info(
            "Waterfall chart cache update: plot-level data label formula %s cached %s points",
            f_node.text,
            len(label_values),
        )
    if label_cache_updates:
        logger.info(
            "Waterfall chart cache update: %s data label caches updated",
            label_cache_updates,
        )
    if c15_label_updates:
        logger.info(
            "Waterfall chart cache update: %s c15 label-range caches updated",
            c15_label_updates,
        )
    elif label_cache_missing:
        logger.info(
            "Waterfall chart cache update: chart is not using value-from-cells labels",
        )
    else:
        logger.info(
            "Waterfall chart cache update: chart is not using value-from-cells labels",
        )


def _update_chart_label_caches(chart, workbook) -> None:
    root = chart.part._element
    nsmap = _chart_namespace_map(root)
    label_refs = root.findall(".//c:dLbls//c:tx//c:strRef", namespaces=nsmap)
    if not label_refs:
        return
    label_cache_updates = 0
    c15_label_updates = 0
    for ref_node in label_refs:
        f_node = ref_node.find("c:f", namespaces=nsmap)
        if f_node is None or not f_node.text:
            logger.info("Chart label cache update: data label ref formula missing")
            continue
        label_ws, label_ref_range, _ = _worksheet_and_range_from_formula(
            workbook, f_node.text
        )
        label_rows = _range_values_from_worksheet(label_ws, label_ref_range)
        label_values = _flatten_cell_values(label_rows)
        if _all_blank(label_values):
            logger.info(
                "Chart label cache update: data label range '%s' is blank",
                f_node.text,
            )
            continue
        str_cache, _ = _ensure_str_cache(ref_node)
        _update_str_cache(
            str_cache,
            ["" if value is None else str(value) for value in label_values],
        )
        label_cache_updates += 1
        d_lbls_node = ref_node
        while d_lbls_node is not None and not d_lbls_node.tag.endswith("dLbls"):
            d_lbls_node = d_lbls_node.getparent()
        if d_lbls_node is not None:
            c15_label_updates += _update_c15_label_range_cache(
                d_lbls_node,
                f_node.text,
                ["" if value is None else str(value) for value in label_values],
                nsmap,
                "chart",
            )
    if label_cache_updates:
        logger.info(
            "Chart label cache update: %s data label caches updated",
            label_cache_updates,
        )
    if c15_label_updates:
        logger.info(
            "Chart label cache update: %s c15 label-range caches updated",
            c15_label_updates,
        )


def _capture_label_columns(ws, series_names: list[str]) -> dict[int, dict[str, list]]:
    label_columns: dict[int, dict[str, list]] = {}
    series_lookup = {str(name).strip().lower() for name in series_names if name}
    for col_idx in range(2, ws.max_column + 1):
        header = ws.cell(row=1, column=col_idx).value
        if not header:
            continue
        header_text = str(header).strip().lower()
        if header_text in series_lookup:
            continue
        values = [
            ws.cell(row=row_idx, column=col_idx).value
            for row_idx in range(2, ws.max_row + 1)
        ]
        label_columns[col_idx] = {"header": header, "values": values}
    return label_columns


def _apply_label_columns(ws, label_columns: dict[int, dict[str, list]], total_rows: int) -> None:
    for col_idx, column in label_columns.items():
        ws.cell(row=1, column=col_idx, value=column["header"])
        values = column["values"]
        if len(values) < total_rows:
            values = values + [None] * (total_rows - len(values))
        for row_offset in range(total_rows):
            ws.cell(row=row_offset + 2, column=col_idx, value=values[row_offset])


def _apply_gathered_waterfall_labels(
    ws,
    label_values: dict[str, list],
    total_rows: int,
) -> set[str]:
    applied_headers: set[str] = set()
    if not label_values:
        return applied_headers
    for header, values in label_values.items():
        col_idx = _find_header_column(ws, [header])
        if col_idx is None:
            continue
        applied_headers.add(_normalize_column_name(header))
        aligned_values = _align_label_values(values, total_rows)
        for row_offset in range(total_rows):
            ws.cell(row=row_offset + 2, column=col_idx, value=aligned_values[row_offset])
    return applied_headers


def _update_lab_base_label(
    label_columns: dict[int, dict[str, list]],
    base_indices: tuple[int, int] | None,
    base_values: tuple[float, float] | None,
    total_rows: int,
) -> None:
    if base_indices is None or base_values is None:
        return
    formatted_values = [
        _format_lab_base_value(value) for value in base_values
    ]
    base_rows = list(base_indices)
    for column in label_columns.values():
        header = str(column["header"]).strip().lower()
        if header != "labs-base":
            continue
        values = column["values"]
        if len(values) < total_rows:
            values.extend([None] * (total_rows - len(values)))
        for idx, base_row in enumerate(base_rows):
            if base_row is None or base_row < 0:
                continue
            if base_row < len(values):
                values[base_row] = formatted_values[idx]
        column["values"] = values
        return


def _is_blank_cell(value) -> bool:
    if value is None:
        return True
    if isinstance(value, str) and not value.strip():
        return True
    return False


def _normalize_header_value(value: str) -> str:
    return str(value).strip().lower()


def _ensure_negatives_column_positive(ws) -> None:
    header_row = None
    header_col = None
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
        for cell in row:
            value = cell.value
            if value is None:
                continue
            if _normalize_header_value(value) == "negatives":
                header_row = cell.row
                header_col = cell.column
                break
        if header_row is not None:
            break
    if header_row is None or header_col is None:
        return

    label_col = header_col - 1 if header_col > 1 else header_col
    empty_streak = 0
    for row_idx in range(header_row + 1, ws.max_row + 1):
        label_value = ws.cell(row=row_idx, column=label_col).value
        negatives_cell = ws.cell(row=row_idx, column=header_col)
        if _is_blank_cell(label_value):
            if _is_blank_cell(negatives_cell.value):
                empty_streak += 1
                if empty_streak >= 2:
                    break
                continue
            break
        empty_streak = 0
        value = negatives_cell.value
        if isinstance(value, numbers.Number) and not isinstance(value, bool):
            negatives_cell.value = abs(value)
            continue
        if isinstance(value, str) and value.lstrip().startswith("="):
            formula = value.lstrip()[1:].strip()
            if not (formula.lower().startswith("abs(") and formula.endswith(")")):
                negatives_cell.value = f"=ABS({formula})"


def _find_header_column(ws, candidates: list[str]) -> int | None:
    normalized_columns = {}
    for col_idx in range(1, ws.max_column + 1):
        value = ws.cell(row=1, column=col_idx).value
        if value is None:
            continue
        normalized_columns[_normalize_column_name(str(value))] = col_idx
    candidate_normalized = [_normalize_column_name(candidate) for candidate in candidates]
    for candidate in candidate_normalized:
        if candidate in normalized_columns:
            return normalized_columns[candidate]
    for column_key, col_idx in normalized_columns.items():
        for candidate in candidate_normalized:
            if candidate in column_key or column_key in candidate:
                return col_idx
    from difflib import get_close_matches

    matches = get_close_matches(
        " ".join(candidate_normalized),
        list(normalized_columns.keys()),
        n=1,
        cutoff=0.75,
    )
    if matches:
        return normalized_columns[matches[0]]
    return None


def _numeric_cell_value(cell) -> float | None:
    value = cell.value
    if value is None:
        return None
    if isinstance(value, numbers.Number) and not isinstance(value, bool):
        return float(value)
    if isinstance(value, str):
        stripped = value.strip()
        if stripped.startswith("="):
            cached_value = cell.internal_value
            if isinstance(cached_value, numbers.Number) and not isinstance(cached_value, bool):
                return float(cached_value)
            cached_value = getattr(cell, "_value", None)
            if isinstance(cached_value, numbers.Number) and not isinstance(cached_value, bool):
                return float(cached_value)
            return None
        try:
            return float(stripped)
        except ValueError:
            return None
    return None


def _format_waterfall_label(value: float, sign: str) -> str:
    abs_value = abs(value)
    if abs_value >= 1_000_000:
        scaled = abs_value / 1_000_000
        suffix = "m"
    elif abs_value >= 1_000:
        scaled = abs_value / 1_000
        suffix = "k"
    else:
        scaled = abs_value
        suffix = ""
    return f"{sign}{scaled:.1f}{suffix}"


def _waterfall_labs_header_for_series(series_name: str | None) -> str | None:
    if not series_name:
        return None
    series_label = str(series_name).strip().lower()
    if "base" in series_label:
        return "labs-Base"
    if "promo" in series_label:
        return "labs-Promo"
    if "media" in series_label:
        return "labs-Media"
    if "blank" in series_label:
        return "labs-Blanks"
    if "positive" in series_label:
        return "labs-Positives"
    if "negative" in series_label:
        return "labs-Negatives"
    return None


def _resolve_waterfall_labs_column(ws, series_name: str | None) -> tuple[str | None, int | None]:
    header = _waterfall_labs_header_for_series(series_name)
    if not header:
        return None, None
    return header, _find_header_column(ws, [header])


def _update_all_waterfall_labs(
    ws,
    base_indices: tuple[int, int] | None,
    base_values: tuple[float, float] | None,
    skip_headers: set[str] | None = None,
) -> None:
    skip_headers = {value for value in (skip_headers or set()) if value}
    labs_base_col = _find_header_column(ws, ["labs-Base"])
    labs_promo_col = _find_header_column(ws, ["labs-Promo"])
    labs_media_col = _find_header_column(ws, ["labs-Media"])
    labs_blanks_col = _find_header_column(ws, ["labs-Blanks"])
    labs_positives_col = _find_header_column(ws, ["labs-Positives"])
    labs_negatives_col = _find_header_column(ws, ["labs-Negatives"])

    promo_col = _find_header_column(ws, ["Promo"])
    media_col = _find_header_column(ws, ["Media"])
    positives_col = _find_header_column(ws, ["Positives"])
    negatives_col = _find_header_column(ws, ["Negatives"])

    total_rows = ws.max_row
    if labs_base_col and _normalize_column_name("labs-Base") not in skip_headers:
        for row_idx in range(2, total_rows + 1):
            ws.cell(row=row_idx, column=labs_base_col).value = None
        if base_indices and base_values:
            formatted = [_format_lab_base_value(value) for value in base_values]
            for idx, base_row in enumerate(base_indices):
                if base_row is None or base_row < 0:
                    continue
                row_idx = base_row + 2
                if row_idx <= total_rows:
                    ws.cell(row=row_idx, column=labs_base_col).value = formatted[idx]

    for row_idx in range(2, total_rows + 1):
        if labs_promo_col and _normalize_column_name("labs-Promo") not in skip_headers:
            value = _numeric_cell_value(ws.cell(row=row_idx, column=promo_col)) if promo_col else None
            cell = ws.cell(row=row_idx, column=labs_promo_col)
            if value is None or value == 0:
                cell.value = None
            else:
                sign = "+" if value > 0 else "-"
                cell.value = _format_waterfall_label(value, sign)
        if labs_media_col and _normalize_column_name("labs-Media") not in skip_headers:
            value = _numeric_cell_value(ws.cell(row=row_idx, column=media_col)) if media_col else None
            cell = ws.cell(row=row_idx, column=labs_media_col)
            if value is None or value == 0:
                cell.value = None
            else:
                sign = "+" if value > 0 else "-"
                cell.value = _format_waterfall_label(value, sign)
        if labs_blanks_col and _normalize_column_name("labs-Blanks") not in skip_headers:
            ws.cell(row=row_idx, column=labs_blanks_col).value = None
        if labs_positives_col and _normalize_column_name("labs-Positives") not in skip_headers:
            pos_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=positives_col))
                if positives_col
                else None
            )
            cell = ws.cell(row=row_idx, column=labs_positives_col)
            if pos_value is None or pos_value == 0:
                cell.value = None
            else:
                cell.value = _format_waterfall_label(pos_value, "+")
        if labs_negatives_col and _normalize_column_name("labs-Negatives") not in skip_headers:
            neg_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=negatives_col))
                if negatives_col
                else None
            )
            cell = ws.cell(row=row_idx, column=labs_negatives_col)
            if neg_value is None or neg_value == 0:
                cell.value = None
            else:
                cell.value = _format_waterfall_label(neg_value, "-")
    logger.info(
        "Waterfall chart labels: overwrote labs columns with literal values (Base=%s, Promo=%s, Media=%s, Blanks=%s, Positives=%s, Negatives=%s).",
        bool(labs_base_col),
        bool(labs_promo_col),
        bool(labs_media_col),
        bool(labs_blanks_col),
        bool(labs_positives_col),
        bool(labs_negatives_col),
    )


def _update_waterfall_positive_negative_labels(ws) -> None:
    positives_col = _find_header_column(ws, ["Positives"])
    negatives_col = _find_header_column(ws, ["Negatives"])
    labs_positives_col = _find_header_column(ws, ["labs-Positives"])
    labs_negatives_col = _find_header_column(ws, ["labs-Negatives"])
    if not positives_col and not negatives_col:
        return
    if not labs_positives_col and not labs_negatives_col:
        return

    for row_idx in range(2, ws.max_row + 1):
        if labs_positives_col:
            pos_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=positives_col))
                if positives_col
                else None
            )
            labs_cell = ws.cell(row=row_idx, column=labs_positives_col)
            if pos_value is None or pos_value == 0:
                labs_cell.value = None
            else:
                labs_cell.value = _format_waterfall_label(pos_value, "+")
        if labs_negatives_col:
            neg_value = (
                _numeric_cell_value(ws.cell(row=row_idx, column=negatives_col))
                if negatives_col
                else None
            )
            labs_cell = ws.cell(row=row_idx, column=labs_negatives_col)
            if neg_value is None or neg_value == 0:
                labs_cell.value = None
            else:
                labs_cell.value = _format_waterfall_label(neg_value, "-")


def _format_yoy_change_text(value: float) -> str:
    if value is None or pd.isna(value):
        return "0%"
    return f"{abs(value):.0%}"


def _remove_shapes_with_text(slide, targets: list[str]) -> None:
    if not targets:
        return
    for shape in list(slide.shapes):
        if not shape.has_text_frame:
            continue
        text_value = shape.text_frame.text
        if any(target in text_value for target in targets):
            element = shape._element
            element.getparent().remove(element)


def _update_waterfall_yoy_arrows(
    slide,
    base_values: tuple[float, float] | None,
) -> None:
    if base_values is None:
        return
    year1_total, year2_total = base_values
    if year1_total is None or year2_total is None:
        return
    if year1_total == 0:
        pct_change = 0.0
    else:
        pct_change = (year2_total - year1_total) / year1_total
    direction = "increase" if year2_total >= year1_total else "decrease"
    remove_placeholder = "<% decrease>" if direction == "increase" else "<% increase>"
    keep_placeholder = "<% increase>" if direction == "increase" else "<% decrease>"
    _remove_shapes_with_text(slide, [remove_placeholder])
    replacement_text = f"{_format_yoy_change_text(pct_change)} {direction}"
    replace_text_in_slide_preserve_formatting(slide, keep_placeholder, replacement_text)


def _set_waterfall_chart_title(chart, label: str | None) -> None:
    if not label:
        return
    title_text = f"{label} Waterfall"
    try:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text
    except Exception:
        return

def update_or_add_column_chart(slide, chart_name, categories, series_dict):
    """
    If a chart with name=chart_name exists on the slide, update its data.
    Otherwise insert a new clustered column chart in a sensible spot.
    Charts produced here remain EDITABLE in PowerPoint.
    """
    chart_shape = None
    for shape in slide.shapes:
        if getattr(shape, "name", None) == chart_name:
            if shape.has_chart:
                chart_shape = shape
                break
            else:
                # Remove placeholder artifacts that aren't real charts
                sp = shape._element
                sp.getparent().remove(sp)

    cd = ChartData()
    cd.categories = categories
    for s_name, values in series_dict.items():
        cd.add_series(s_name, list(values))

    if chart_shape:
        # Replace data in existing chart (preserves template styling)
        chart_shape.chart.replace_data(cd)
        updated_wb = _load_chart_workbook(chart_shape.chart)
        _update_chart_label_caches(chart_shape.chart, updated_wb)
        return chart_shape
    else:
        # Fallback: repurpose the first chart on the slide if present.
        for shape in slide.shapes:
            if shape.has_chart:
                shape.chart.replace_data(cd)
                shape.name = chart_name
                updated_wb = _load_chart_workbook(shape.chart)
                _update_chart_label_caches(shape.chart, updated_wb)
                return shape
        # Remove any stale shapes with the target name before adding a new chart
        for shape in list(slide.shapes):
            if getattr(shape, "name", None) == chart_name:
                sp = shape._element
                sp.getparent().remove(sp)

        # Insert a new chart (fallback)
        left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, cd
        )
        chart_shape.name = chart_name
        chart = chart_shape.chart
        # Light touch formatting; rely on template/theme for styling
        chart.has_legend = True
        updated_wb = _load_chart_workbook(chart)
        _update_chart_label_caches(chart, updated_wb)
        return chart

def update_or_add_waterfall_chart(slide, chart_name, categories, series_dict):
    """
    Update an existing waterfall chart or insert one if missing.
    """
    chart_shape = None
    for shape in slide.shapes:
        if getattr(shape, "name", None) == chart_name and shape.has_chart:
            chart_shape = shape
            break

    if chart_shape is None:
        for shape in slide.shapes:
            if shape.has_chart:
                chart_shape = shape
                shape.name = chart_name
                break

    cd = ChartData()
    cd.categories = categories
    for s_name, values in series_dict.items():
        cd.add_series(s_name, list(values))

    if chart_shape:
        chart_shape.chart.replace_data(cd)
        return chart_shape

    waterfall_type = getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED)
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4.5)
    chart_shape = slide.shapes.add_chart(
        waterfall_type, left, top, width, height, cd
    )
    chart_shape.name = chart_name
    return chart_shape

def set_text_by_name(slide, shape_name, text):
    for shape in slide.shapes:
        if getattr(shape, "name", None) == shape_name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(text)
            p.alignment = PP_ALIGN.LEFT
            return True
    return False


def replace_text_in_slide(slide, old_text, new_text):
    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        current_text = text_frame.text
        if old_text in current_text:
            text_frame.text = current_text.replace(old_text, new_text)
            replaced = True
    return replaced


def replace_text_in_slide_preserve_formatting(slide, old_text, new_text):
    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            found_in_runs = False
            for run in paragraph.runs:
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    found_in_runs = True
                    replaced = True
            if found_in_runs or old_text not in paragraph.text:
                continue
            updated_text = paragraph.text.replace(old_text, new_text)
            _rebuild_paragraph_runs(paragraph, updated_text)
            replaced = True
    return replaced


def _replace_placeholder_in_paragraph_runs(paragraph, placeholder: str, replacement: str) -> bool:
    if not replacement:
        return False
    runs = list(paragraph.runs)
    if not runs:
        return False
    replaced = False
    while True:
        full_text = "".join(run.text for run in runs)
        start_idx = full_text.find(placeholder)
        if start_idx == -1:
            break
        end_idx = start_idx + len(placeholder)
        replaced = True
        first_overlap = True
        cumulative = 0
        for run in runs:
            run_text = run.text
            run_start = cumulative
            run_end = cumulative + len(run_text)
            cumulative = run_end
            if run_end <= start_idx or run_start >= end_idx:
                continue
            overlap_start = max(start_idx, run_start)
            overlap_end = min(end_idx, run_end)
            local_start = overlap_start - run_start
            local_end = overlap_end - run_start
            if first_overlap:
                run.text = run_text[:local_start] + replacement + run_text[local_end:]
                first_overlap = False
            else:
                run.text = run_text[:local_start] + run_text[local_end:]
    return replaced


def _replace_placeholders_in_slide_runs(
    slide, replacements: dict[str, str | None]
) -> bool:
    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for placeholder, value in replacements.items():
                if not value:
                    continue
                if _replace_placeholder_in_paragraph_runs(paragraph, placeholder, value):
                    replaced = True
    return replaced


def _capture_run_formatting(run):
    font = run.font
    color = font.color
    return {
        "name": font.name,
        "size": font.size,
        "bold": font.bold,
        "italic": font.italic,
        "underline": font.underline,
        "color_rgb": color.rgb if color is not None else None,
    }


def _apply_run_formatting(run, formatting):
    font = run.font
    font.name = formatting["name"]
    font.size = formatting["size"]
    font.bold = formatting["bold"]
    font.italic = formatting["italic"]
    font.underline = formatting["underline"]
    if formatting["color_rgb"] is not None:
        font.color.rgb = formatting["color_rgb"]


def _rebuild_paragraph_runs(paragraph, new_text: str) -> None:
    original_runs = list(paragraph.runs)
    if not original_runs:
        paragraph.text = new_text
        return
    formats = [_capture_run_formatting(run) for run in original_runs]
    run_lengths = [len(run.text) for run in original_runs]
    for run in original_runs:
        paragraph._element.remove(run._r)
    cursor = 0
    for idx, fmt in enumerate(formats):
        if idx == len(formats) - 1:
            segment = new_text[cursor:]
        else:
            segment = new_text[cursor : cursor + run_lengths[idx]]
        new_run = paragraph.add_run()
        new_run.text = segment
        _apply_run_formatting(new_run, fmt)
        cursor += len(segment)


def _shape_text_snippet(shape) -> str:
    if not shape.has_text_frame:
        return ""
    text = shape.text_frame.text or ""
    compact = " ".join(text.split())
    return compact[:80]


def _chart_title_text_frame(chart):
    try:
        if chart.has_title:
            return chart.chart_title.text_frame
    except Exception:
        return None
    return None


def _slide_title(slide) -> str:
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None
    if title_shape is not None and title_shape.has_text_frame:
        title_text = title_shape.text_frame.text or ""
        return title_text.strip()
    return ""


def _slide_index(prs, target_slide) -> int | None:
    for idx, slide in enumerate(prs.slides, start=1):
        if slide is target_slide:
            return idx
    return None


def _slides_with_placeholder(prs, placeholder: str) -> list[int]:
    matches: list[int] = []
    for idx, slide in enumerate(prs.slides, start=1):
        found = False
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder in (shape.text_frame.text or ""):
                found = True
                break
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if placeholder in (cell.text_frame.text or ""):
                            found = True
                            break
                    if found:
                        break
            if found:
                break
            if shape.has_chart:
                chart_text_frame = _chart_title_text_frame(shape.chart)
                if chart_text_frame and placeholder in (chart_text_frame.text or ""):
                    found = True
                    break
        if found:
            matches.append(idx)
    return matches


def _replace_placeholders_in_text_frame(text_frame, replacements, counts) -> None:
    for paragraph in text_frame.paragraphs:
        for placeholder, replacement in replacements.items():
            paragraph_text = paragraph.text or ""
            occurrences = paragraph_text.count(placeholder)
            if occurrences == 0:
                continue
            counts[placeholder]["found"] += occurrences
            if replacement is None:
                continue
            if paragraph.runs:
                if _replace_placeholder_in_paragraph_runs(paragraph, placeholder, replacement):
                    counts[placeholder]["replaced"] += occurrences
            else:
                paragraph.text = paragraph_text.replace(placeholder, replacement)
                counts[placeholder]["replaced"] += occurrences


def replace_placeholders_strict(prs, slide_selector, replacements: dict[str, str | None]) -> None:
    if slide_selector is None:
        raise ValueError("Slide selector is required to replace placeholders.")
    if isinstance(slide_selector, str):
        slide = _find_slide_by_marker(prs, slide_selector)
    else:
        slide = slide_selector
    if slide is None:
        raise ValueError(f"Could not resolve slide selector: {slide_selector}")

    counts = {
        placeholder: {"found": 0, "replaced": 0} for placeholder in replacements
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            _replace_placeholders_in_text_frame(shape.text_frame, replacements, counts)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_placeholders_in_text_frame(
                        cell.text_frame, replacements, counts
                    )
        if shape.has_chart:
            chart_text_frame = _chart_title_text_frame(shape.chart)
            if chart_text_frame is not None:
                _replace_placeholders_in_text_frame(
                    chart_text_frame, replacements, counts
                )

    slide_idx = _slide_index(prs, slide)
    slide_title = _slide_title(slide)
    slide_name = getattr(slide, "name", None) or ""
    shape_lines = []
    for shape in slide.shapes:
        shape_lines.append(
            " - "
            f"id={getattr(shape, 'shape_id', None)} "
            f"name={getattr(shape, 'name', None)!r} "
            f"type={getattr(shape, 'shape_type', None)} "
            f"has_text_frame={shape.has_text_frame} "
            f"has_table={shape.has_table} "
            f"text={_shape_text_snippet(shape)!r}"
        )
        if shape.has_chart:
            chart_text_frame = _chart_title_text_frame(shape.chart)
            if chart_text_frame is not None:
                chart_text = chart_text_frame.text or ""
                compact = " ".join(chart_text.split())
                shape_lines.append(
                    " - "
                    f"chart_title={compact[:80]!r} "
                    f"shape_name={getattr(shape, 'name', None)!r}"
                )
    counts_lines = [
        f" - {placeholder}: found={stats['found']} replaced={stats['replaced']}"
        for placeholder, stats in counts.items()
    ]

    def build_diagnostics(missing_placeholder: str) -> str:
        locations = _slides_with_placeholder(prs, missing_placeholder)
        location_line = (
            f"Slides containing {missing_placeholder}: {locations}"
            if locations
            else f"Slides containing {missing_placeholder}: []"
        )
        return "\n".join(
            [
                "Placeholder replacement diagnostics:",
                f"Slide index: {slide_idx}",
                f"Slide name: {slide_name}",
                f"Slide title: {slide_title}",
                "Shape inventory:",
                *shape_lines,
                "Replacement counts:",
                *counts_lines,
                location_line,
            ]
        )

    for placeholder in replacements:
        if counts[placeholder]["found"] > 0:
            continue
        locations = _slides_with_placeholder(prs, placeholder)
        if not locations:
            raise ValueError(
                f"Placeholder {placeholder} not found in deck\n"
                f"{build_diagnostics(placeholder)}"
            )
        intended_idx = slide_idx if slide_idx is not None else "unknown"
        raise ValueError(
            f"Placeholder {placeholder} found on slide {locations[0]} "
            f"not on Waterfall slide {intended_idx}\n"
            f"{build_diagnostics(placeholder)}"
        )


def _update_waterfall_axis_placeholders(
    prs,
    slide_selector,
    target_level_label_value: str | None,
    modelled_in_value: str | None,
    metric_value: str | None,
) -> None:
    replacements = {
        "<Target Level Label>": target_level_label_value,
        "<modelled in>": modelled_in_value,
        "<metric>": metric_value,
    }
    replace_placeholders_strict(prs, slide_selector, replacements)
    if not modelled_in_value:
        logger.warning(
            "Missing/blank value for 'Sales will be modelled in:' in Project Details."
        )
    if not metric_value:
        logger.warning(
            "Missing/blank value for 'Volume metric (unique per dataset):' in Project Details."
        )


def append_text_after_label(slide, label_text, appended_text):
    if not appended_text:
        return False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if label_text not in paragraph.text:
                continue
            if appended_text in paragraph.text:
                return True
            spacer = "" if label_text.endswith(" ") else " "
            for run in paragraph.runs:
                if label_text in run.text:
                    new_run = paragraph.add_run()
                    new_run.text = f"{spacer}{appended_text}"
                    return True
            new_run = paragraph.add_run()
            new_run.text = f"{spacer}{appended_text}"
            return True
    return False


def append_paragraph_after_label(slide, label_text, appended_text):
    if not appended_text:
        return False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        if label_text not in text_frame.text:
            continue
        if any(paragraph.text.strip() == appended_text for paragraph in text_frame.paragraphs):
            return True
        for paragraph in text_frame.paragraphs:
            if label_text in paragraph.text:
                new_paragraph = text_frame.add_paragraph()
                new_paragraph.text = appended_text
                paragraph._p.addnext(new_paragraph._p)
                return True
    return False


def append_paragraphs_after_label(slide, label_text, appended_texts):
    if not appended_texts:
        return False
    appended_texts = [text for text in appended_texts if text and text.strip()]
    if not appended_texts:
        return False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        if label_text not in text_frame.text:
            continue
        existing_texts = {paragraph.text.strip() for paragraph in text_frame.paragraphs}
        to_add = [text for text in appended_texts if text not in existing_texts]
        if not to_add:
            return True
        insert_after = None
        for paragraph in text_frame.paragraphs:
            if label_text in paragraph.text:
                insert_after = paragraph
                break
        if insert_after is None:
            continue
        last_paragraph = insert_after
        for text in to_add:
            new_paragraph = text_frame.add_paragraph()
            new_paragraph.text = text
            last_paragraph._p.addnext(new_paragraph._p)
            last_paragraph = new_paragraph
        return True
    return False

def add_table(slide, table_name, df: pd.DataFrame):
    # Identify an existing table to reuse, preferring one with the expected name.
    target_shape = None
    for shape in slide.shapes:
        if getattr(shape, "name", None) == table_name and shape.has_table:
            target_shape = shape
            break

    if target_shape is None:
        for shape in slide.shapes:
            if shape.has_table:
                target_shape = shape
                target_shape.name = table_name
                break

    if target_shape and target_shape.has_table:
        tbl = target_shape.table
        # Resize (simple): write headers to row 0, then rows afterward if room allows
        n_rows = min(len(df) + 1, tbl.rows.__len__())
        n_cols = min(len(df.columns), tbl.columns.__len__())
        # headers
        for j, col in enumerate(df.columns[:n_cols]):
            cell = tbl.cell(0, j)
            cell.text = str(col)
        # cells
        for i in range(1, n_rows):
            for j in range(n_cols):
                tbl.cell(i, j).text = str(df.iloc[i-1, j])
        # Clear any leftover rows beyond the populated range
        for i in range(n_rows, tbl.rows.__len__()):
            for j in range(tbl.columns.__len__()):
                tbl.cell(i, j).text = ""
        return True

    # Remove non-table placeholders with the desired name so we can insert a fresh table.
    for shape in list(slide.shapes):
        if getattr(shape, "name", None) == table_name and not shape.has_table:
            sp = shape._element
            sp.getparent().remove(sp)

    # Otherwise, add a new table
    rows, cols = len(df) + 1, len(df.columns)
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(1 + 0.3 * rows)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table_shape.name = table_name
    table = table_shape.table
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    for i in range(len(df)):
        for j in range(len(df.columns)):
            table.cell(i+1, j).text = str(df.iloc[i, j])

    # Avoid manipulating the low-level XML that may not exist across templates.
    # python-pptx represents the table as a ``CT_GraphicalObjectFrame`` whose
    # schema does not expose a ``graphicFrame`` attribute.  Some versions of the
    # library can therefore raise an AttributeError when we try to clear borders
    # by touching ``graphicFrame`` directly.  Since this styling tweak is only a
    # nice-to-have, we simply rely on the template/theme defaults instead of
    # editing the XML manually.
    return True

def remove_empty_placeholders(slide):
    """Remove placeholder shapes that have no meaningful content."""
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue

        # Keep placeholders that now contain text, tables, or charts with data.
        if shape.has_text_frame:
            if shape.text_frame.text and shape.text_frame.text.strip():
                continue
        elif shape.has_table:
            # If every cell is blank, treat as empty.
            if any(
                cell.text.strip()
                for row in shape.table.rows
                for cell in row.cells
            ):
                continue
        elif shape.has_chart:
            # Assume populated charts should remain.
            continue

        sp = shape._element
        sp.getparent().remove(sp)


def _normalize_label(value: str) -> str:
    return " ".join(value.strip().lower().replace(":", "").split())


def _find_company_week_value(scope_df: pd.DataFrame, label: str):
    if scope_df is None or scope_df.empty or scope_df.shape[1] < 2:
        raise ValueError("Scope file must include labels in column A and company weeks in column B.")
    label_normalized = _normalize_label(label)
    candidates = []
    for _, row in scope_df.iterrows():
        cell_value = row.iloc[0]
        if pd.isna(cell_value):
            continue
        cell_label = _normalize_label(str(cell_value))
        candidates.append((cell_label, row))
        if label_normalized in cell_label or cell_label in label_normalized:
            value = row.iloc[1]
            if pd.isna(value):
                raise ValueError(f"Missing company week value for '{label}'.")
            return value
    if candidates:
        from difflib import get_close_matches

        labels = [candidate_label for candidate_label, _ in candidates]
        matches = get_close_matches(label_normalized, labels, n=1, cutoff=0.7)
        if matches:
            matched_label = matches[0]
            for candidate_label, row in candidates:
                if candidate_label == matched_label:
                    value = row.iloc[1]
                    if pd.isna(value):
                        raise ValueError(f"Missing company week value for '{label}'.")
                    return value
    raise ValueError(f"Could not find '{label}' in the scope file.")


def _coerce_yearwk(value) -> int:
    if pd.isna(value):
        raise ValueError("Missing company week value for modelling period.")
    raw = str(value).strip()
    if not raw:
        raise ValueError("Missing company week value for modelling period.")
    try:
        numeric_value = int(float(raw))
    except ValueError:
        digits = "".join(ch for ch in raw if ch.isdigit())
        if not digits:
            raise ValueError("Company week value must be a YYYYWW-style week number or a company week.")
        numeric_value = int(digits)

    if len(str(numeric_value)) <= 4:
        mapper = CompanyWeekMapper(
            anchor_company_week=2455,
            anchor_yearwk=202638,
            check_company_week=2470,
            check_yearwk=202653,
        )
        return mapper.to_yearwk(numeric_value)

    year, week = divmod(numeric_value, 100)
    if year <= 0 or not (1 <= week <= 53):
        raise ValueError("Company week value must be a YYYYWW-style week number or a company week.")
    return numeric_value


def _format_modelling_period(data_df: pd.DataFrame, scope_df: pd.DataFrame) -> tuple[str, int]:
    start_company_week = _find_company_week_value(scope_df, "First week of modelling")
    end_company_week = _find_company_week_value(scope_df, "Last week of modelling")
    start_yearwk = _coerce_yearwk(start_company_week)
    end_yearwk = _coerce_yearwk(end_company_week)
    start_date = CompanyWeekMapper._yearwk_to_monday(start_yearwk)
    end_date = CompanyWeekMapper._yearwk_to_monday(end_yearwk) + timedelta(days=6)
    week_count = ((end_date - start_date).days // 7) + 1
    return f"{start_date:%b %d, %Y} - {end_date:%b %d, %Y}", week_count


def _format_study_year_range(scope_df: pd.DataFrame) -> str:
    start_company_week = _find_company_week_value(scope_df, "First week of modelling")
    end_company_week = _find_company_week_value(scope_df, "Last week of modelling")
    start_yearwk = _coerce_yearwk(start_company_week)
    end_yearwk = _coerce_yearwk(end_company_week)
    start_year = CompanyWeekMapper._yearwk_to_monday(start_yearwk).year
    end_year = CompanyWeekMapper._yearwk_to_monday(end_yearwk).year
    if start_year == end_year:
        return str(start_year)
    return f"{start_year}-{end_year}"


def set_time_period_text(slide, label_text, time_period, week_count):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
        if label_text not in text_frame.text:
            continue
        label_index = None
        for idx, paragraph in enumerate(text_frame.paragraphs):
            if label_text in paragraph.text:
                label_index = idx
                break
        if label_index is None:
            continue
        paragraphs = list(text_frame.paragraphs)
        if label_index + 1 < len(paragraphs):
            value_paragraph = paragraphs[label_index + 1]
        else:
            value_paragraph = text_frame.add_paragraph()
        value_paragraph.text = f"{time_period} (number of weeks = {week_count})"
        for extra_paragraph in paragraphs[label_index + 2:]:
            extra_paragraph.text = ""
        return True
    return False


def _modelling_period_bounds(scope_df: pd.DataFrame) -> tuple[date, date]:
    start_company_week = _find_company_week_value(scope_df, "First week of modelling")
    end_company_week = _find_company_week_value(scope_df, "Last week of modelling")
    start_yearwk = _coerce_yearwk(start_company_week)
    end_yearwk = _coerce_yearwk(end_company_week)
    earliest_yearwk = min(start_yearwk, end_yearwk)
    latest_yearwk = max(start_yearwk, end_yearwk)
    earliest_date = CompanyWeekMapper._yearwk_to_monday(earliest_yearwk)
    latest_date = CompanyWeekMapper._yearwk_to_monday(latest_yearwk) + timedelta(days=6)
    return earliest_date, latest_date


def _replace_modelling_period_placeholders_in_categories(
    categories: list[str],
    scope_df: pd.DataFrame | None,
) -> list[str]:
    if scope_df is None or not categories:
        return categories
    try:
        start_date, end_date = _modelling_period_bounds(scope_df)
    except Exception:
        return categories
    earliest = start_date.strftime("%b %d, %Y")
    latest = end_date.strftime("%b %d, %Y")
    updated = []
    for value in categories:
        text = "" if value is None else str(value)
        if "<earliest date>" in text or "<latest date>" in text:
            text = text.replace("<earliest date>", earliest)
            text = text.replace("<latest date>", latest)
        updated.append(text)
    return updated


def _waterfall_chart_from_slide(slide, chart_name: str):
    for shape in slide.shapes:
        if shape.has_chart and chart_name in (getattr(shape, "name", "") or ""):
            return shape.chart
    for shape in slide.shapes:
        if shape.has_chart:
            return shape.chart
    return None


def _waterfall_chart_shape_from_slide(slide, chart_name: str):
    for shape in slide.shapes:
        if shape.has_chart and chart_name in (getattr(shape, "name", "") or ""):
            return shape
    for shape in slide.shapes:
        if shape.has_chart:
            return shape
    return None


def _clone_chart_part(chart_part):
    package = chart_part.package
    new_partname = package.next_partname(chart_part.partname_template)
    new_chart_part = chart_part.__class__.load(
        new_partname,
        chart_part.content_type,
        package,
        chart_part.blob,
    )
    xlsx_part = chart_part.chart_workbook.xlsx_part
    if xlsx_part is not None:
        new_xlsx_part = EmbeddedXlsxPart.new(xlsx_part.blob, package)
        new_chart_part.chart_workbook.xlsx_part = new_xlsx_part
    return new_chart_part


def _ensure_unique_chart_parts_on_slide(slide, seen_partnames: set[str]) -> None:
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        chart_part = shape.chart.part
        partname = str(chart_part.partname)
        if partname in seen_partnames:
            new_chart_part = _clone_chart_part(chart_part)
            new_rid = shape.part.relate_to(new_chart_part, RT.CHART)
            chart_element = shape._element.graphic.graphicData.chart
            chart_element.set(qn("r:id"), new_rid)
            chart_part = shape.chart.part
            partname = str(chart_part.partname)
        seen_partnames.add(partname)


def _categories_from_chart(chart) -> list[str]:
    categories = []
    try:
        plot_categories = chart.plots[0].categories
    except Exception:
        plot_categories = []
    for category in plot_categories:
        label = getattr(category, "label", None)
        categories.append(str(label) if label is not None else str(category))
    return categories


def _waterfall_base_indices(categories: list[str]) -> tuple[int, int] | None:
    earliest_idx = None
    latest_idx = None
    for idx, value in enumerate(categories):
        text = "" if value is None else str(value)
        if "<earliest date>" in text:
            earliest_idx = idx
        if "<latest date>" in text:
            latest_idx = idx
    if earliest_idx is None or latest_idx is None:
        matches = [
            idx
            for idx, value in enumerate(categories)
            if "52 w/e" in ("" if value is None else str(value)).lower()
        ]
        if len(matches) >= 2:
            earliest_idx = matches[0] if earliest_idx is None else earliest_idx
            latest_idx = matches[-1] if latest_idx is None else latest_idx
    if earliest_idx is None or latest_idx is None:
        return None
    return earliest_idx, latest_idx


def _apply_bucket_categories(
    categories: list[str],
    bucket_labels: list[str],
    base_indices: tuple[int, int],
) -> tuple[list[str], tuple[int, int]]:
    if not bucket_labels:
        return categories, base_indices
    start_idx, end_idx = base_indices
    if start_idx > end_idx:
        start_idx, end_idx = end_idx, start_idx
    prefix = categories[: start_idx + 1]
    suffix = categories[end_idx:]
    new_categories = prefix + bucket_labels + suffix
    new_end_idx = start_idx + len(bucket_labels) + 1
    return new_categories, (start_idx, new_end_idx)


def _apply_bucket_values(
    values: list[float],
    base_indices: tuple[int, int],
    bucket_values: list[float],
) -> list[float]:
    if not bucket_values:
        return values
    start_idx, end_idx = base_indices
    if start_idx > end_idx:
        start_idx, end_idx = end_idx, start_idx
    if end_idx >= len(values):
        values = values + [0.0] * (end_idx - len(values) + 1)
    prefix = values[: start_idx + 1]
    suffix = values[end_idx:]
    return prefix + bucket_values + suffix


def _apply_bucket_placeholders(
    values: list[float],
    base_indices: tuple[int, int],
    bucket_count: int,
    fill_value: float = 0.0,
) -> list[float]:
    if bucket_count <= 0:
        return values
    fill_values = [fill_value] * bucket_count
    return _apply_bucket_values(values, base_indices, fill_values)


def _should_update_base_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    if not name:
        return False
    return "base" in str(name).lower()


def _is_blanks_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    return "blank" in str(name).lower()


def _is_positive_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    return "positive" in str(name).lower()


def _is_negative_series(chart_series) -> bool:
    name = getattr(chart_series, "name", "")
    return "negative" in str(name).lower()


def _bucket_value_split(bucket_values: list[float]) -> tuple[list[float], list[float]]:
    positives: list[float] = []
    negatives: list[float] = []
    for value in bucket_values:
        if value >= 0:
            positives.append(value)
            negatives.append(0.0)
        else:
            positives.append(0.0)
            negatives.append(value)
    return positives, negatives


def _bucket_blank_values(bucket_values: list[float], base_value: float) -> list[float]:
    blanks: list[float] = []
    running_total = base_value
    for value in bucket_values:
        blanks.append(running_total)
        running_total += value
    return blanks


def _align_series_values(values: list[float], total_count: int) -> list[float]:
    if total_count <= 0:
        return values
    if len(values) < total_count:
        return values + [0.0] * (total_count - len(values))
    if len(values) > total_count:
        return values[:total_count]
    return values


def _align_label_values(values: list, total_count: int) -> list:
    if total_count <= 0:
        return values
    if len(values) < total_count:
        return values + [None] * (total_count - len(values))
    if len(values) > total_count:
        return values[:total_count]
    return values


def _sanitize_numeric_value(
    value,
    *,
    label: str | None,
    field: str,
    bucket: str | None = None,
    year1: str | None = None,
    year2: str | None = None,
) -> float:
    if value is None or pd.isna(value):
        logger.warning(
            '[waterfall][sanitize] label="%s" bucket="%s" year1="%s" year2="%s" field="%s" was=%r -> 0.0',
            label or "",
            bucket or "",
            year1 or "",
            year2 or "",
            field,
            value,
        )
        return 0.0
    if isinstance(value, str):
        try:
            return float(value)
        except ValueError as exc:
            raise ValueError(
                "Non-numeric value for "
                f'{field} (label="{label}", bucket="{bucket}", year1="{year1}", year2="{year2}"): {value!r}'
            ) from exc
    try:
        return float(value)
    except (TypeError, ValueError) as exc:
        raise ValueError(
            "Non-numeric value for "
            f'{field} (label="{label}", bucket="{bucket}", year1="{year1}", year2="{year2}"): {value!r}'
        ) from exc


def _sanitize_numeric_list(
    values: list,
    *,
    label: str | None,
    field_prefix: str,
    categories: list[str] | None = None,
    bucket_labels: list[str] | None = None,
    year1: str | None = None,
    year2: str | None = None,
) -> list[float]:
    sanitized: list[float] = []
    for idx, value in enumerate(values):
        bucket = None
        if bucket_labels and idx < len(bucket_labels):
            bucket = bucket_labels[idx]
        elif categories and idx < len(categories):
            bucket = categories[idx]
        sanitized.append(
            _sanitize_numeric_value(
                value,
                label=label,
                field=f"{field_prefix}[{idx}]",
                bucket=bucket,
                year1=year1,
                year2=year2,
            )
        )
    return sanitized


def _waterfall_series_from_gathered_df(
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None,
    target_level_label: str,
) -> tuple[list[str], dict[str, list[float]], dict[str, list]] | None:
    waterfall_df = _build_category_waterfall_df(
        gathered_df,
        target_level_label=target_level_label,
    )
    if waterfall_df.empty:
        return None
    categories = (
        waterfall_df["Vars"]
        .fillna("")
        .astype(str)
        .tolist()
    )
    categories = _replace_modelling_period_placeholders_in_categories(categories, scope_df)
    series_values = {}
    for key in ["Base", "Promo", "Media", "Blanks", "Positives", "Negatives"]:
        if key in waterfall_df.columns:
            series_values[key] = (
                pd.to_numeric(waterfall_df[key], errors="coerce")
                .fillna(0)
                .astype(float)
                .tolist()
            )
    label_values: dict[str, list] = {}
    for key in [
        "labs-Base",
        "labs-Promo",
        "labs-Media",
        "labs-Blanks",
        "labs-Positives",
        "labs-Negatives",
    ]:
        if key in waterfall_df.columns:
            values = []
            for value in waterfall_df[key].tolist():
                if pd.isna(value):
                    values.append(None)
                else:
                    values.append(value)
            label_values[key] = values
    if not series_values and not label_values:
        return None
    return categories, series_values, label_values


def _build_waterfall_chart_data(
    chart,
    scope_df: pd.DataFrame | None,
    gathered_df: pd.DataFrame | None = None,
    target_level_label: str | None = None,
    bucket_labels: list[str] | None = None,
    bucket_values: list[float] | None = None,
    year1: str | None = None,
    year2: str | None = None,
) -> tuple[
    ChartData,
    list[str],
    tuple[int, int] | None,
    tuple[float, float] | None,
    list[tuple[str, list[float]]],
    dict[str, list],
]:
    gathered_override = None
    gathered_label_values: dict[str, list] = {}
    if gathered_df is not None and target_level_label:
        try:
            gathered_override = _waterfall_series_from_gathered_df(
                gathered_df,
                scope_df,
                target_level_label,
            )
        except Exception as exc:
            logger.info(
                "Skipping gatheredCN10 waterfall data for %r: %s",
                target_level_label,
                exc,
            )
            gathered_override = None
    if gathered_override is not None:
        categories, gathered_series, gathered_label_values = gathered_override
    else:
        categories = _categories_from_chart(chart)
        gathered_series = {}
        categories = _replace_modelling_period_placeholders_in_categories(categories, scope_df)
        gathered_label_values = {}
    base_indices = _waterfall_base_indices(categories)
    original_base_indices = base_indices
    bucket_labels = list(bucket_labels or [])
    bucket_values = _sanitize_numeric_list(
        list(bucket_values or []),
        label=target_level_label,
        field_prefix="bucket_values",
        bucket_labels=bucket_labels,
        year1=year1,
        year2=year2,
    )
    if bucket_labels and bucket_values:
        bucket_len = min(len(bucket_labels), len(bucket_values))
        bucket_labels = bucket_labels[:bucket_len]
        bucket_values = bucket_values[:bucket_len]
    if bucket_labels and base_indices:
        categories, base_indices = _apply_bucket_categories(
            categories,
            bucket_labels,
            base_indices,
        )
    bucket_count = len(bucket_labels)
    base_values = None
    if (
        gathered_df is not None
        and target_level_label
        and base_indices is not None
    ):
        base_values = _waterfall_base_values(
            gathered_df,
            target_level_label,
            year1=year1,
            year2=year2,
        )
        base_values = (
            _sanitize_numeric_value(
                base_values[0],
                label=target_level_label,
                field="base_values[0]",
                year1=year1,
                year2=year2,
            ),
            _sanitize_numeric_value(
                base_values[1],
                label=target_level_label,
                field="base_values[1]",
                year1=year1,
                year2=year2,
            ),
        )
    cd = ChartData()
    cd.categories = categories
    base_start_value = None
    if base_values and base_values[0] is not None:
        base_start_value = float(base_values[0])
    elif base_indices is not None:
        for series in chart.series:
            if _should_update_base_series(series):
                series_values = list(series.values)
                if base_indices[0] < len(series_values):
                    base_start_value = float(series_values[base_indices[0]])
                break
    if base_start_value is None:
        base_start_value = 0.0

    positive_bucket_values = []
    negative_bucket_values = []
    blank_bucket_values = []
    if bucket_labels and bucket_values:
        positive_bucket_values, negative_bucket_values = _bucket_value_split(bucket_values)
        blank_bucket_values = _bucket_blank_values(bucket_values, base_start_value)

    series_candidates = list(gathered_series.keys())
    series_values: list[tuple[str, list[float]]] = []
    for series in chart.series:
        values = list(series.values)
        if gathered_series:
            resolved_series = None
            try:
                resolved_series = _resolve_label_from_text(
                    str(series.name),
                    series_candidates,
                )
            except Exception as exc:
                logger.info(
                    "No gatheredCN10 series match for %r: %s",
                    series.name,
                    exc,
                )
            if resolved_series:
                values = list(gathered_series.get(resolved_series, values))
                values = _align_series_values(values, len(categories))
        if original_base_indices and bucket_labels:
            if _is_positive_series(series):
                if positive_bucket_values:
                    values = _apply_bucket_values(
                        values,
                        original_base_indices,
                        positive_bucket_values,
                    )
                else:
                    values = _apply_bucket_placeholders(
                        values,
                        original_base_indices,
                        bucket_count,
                    )
            elif _is_negative_series(series):
                if negative_bucket_values:
                    values = _apply_bucket_values(
                        values,
                        original_base_indices,
                        negative_bucket_values,
                    )
                else:
                    values = _apply_bucket_placeholders(
                        values,
                        original_base_indices,
                        bucket_count,
                    )
            elif _is_blanks_series(series):
                if blank_bucket_values:
                    values = _apply_bucket_values(
                        values,
                        original_base_indices,
                        blank_bucket_values,
                    )
                else:
                    values = _apply_bucket_placeholders(
                        values,
                        original_base_indices,
                        bucket_count,
                    )
            else:
                values = _apply_bucket_placeholders(
                    values,
                    original_base_indices,
                    bucket_count,
                )
        if base_values and base_indices:
            should_update = _should_update_base_series(series)
            if not should_update and len(chart.series) == 1:
                should_update = True
            if should_update:
                if base_indices[0] < len(values):
                    values[base_indices[0]] = base_values[0]
                if base_indices[1] < len(values):
                    values[base_indices[1]] = base_values[1]
        values = _sanitize_numeric_list(
            values,
            label=target_level_label,
            field_prefix=f"series_values[{len(series_values)}]",
            categories=categories,
            year1=year1,
            year2=year2,
        )
        cd.add_series(series.name, values)
        series_values.append((series.name, values))
    return cd, categories, base_indices, base_values, series_values, gathered_label_values


@dataclass(frozen=True)
class WaterfallPayload:
    categories: list[str]
    series_values: list[tuple[str, list[float]]]
    base_indices: tuple[int, int] | None
    base_values: tuple[float, float] | None
    gathered_label_values: dict[str, list]


def _payload_checksum(series_values: list[tuple[str, list[float]]]) -> float:
    if not series_values:
        return 0.0
    checksum = 0.0
    if isinstance(series_values[0], tuple):
        for series_idx, (_, values) in enumerate(series_values):
            for value_idx, value in enumerate(values):
                if value is None or pd.isna(value):
                    logger.warning(
                        '[waterfall][checksum] field="series_values[%d][%d]" was=%r -> 0.0',
                        series_idx,
                        value_idx,
                        value,
                    )
                    continue
                checksum += abs(float(value))
        return checksum
    for value_idx, value in enumerate(series_values):
        if value is None or pd.isna(value):
            logger.warning(
                '[waterfall][checksum] field="values[%d]" was=%r -> 0.0',
                value_idx,
                value,
            )
            continue
        checksum += abs(float(value))
    return checksum


def _chart_data_from_payload(payload: WaterfallPayload) -> ChartData:
    cd = ChartData()
    cd.categories = payload.categories
    for name, values in payload.series_values:
        cd.add_series(name, values)
    return cd


def compute_payload_for_label(
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None,
    target_level_label: str,
    bucket_data: dict | None,
    template_chart,
) -> WaterfallPayload:
    if template_chart is None:
        raise ValueError("Template chart is required to compute waterfall payloads.")
    (
        _,
        categories,
        base_indices,
        base_values,
        series_values,
        gathered_label_values,
    ) = _build_waterfall_chart_data(
        template_chart,
        scope_df,
        gathered_df,
        target_level_label,
        bucket_data.get("labels") if bucket_data else None,
        bucket_data.get("values") if bucket_data else None,
        year1=bucket_data.get("year1") if bucket_data else None,
        year2=bucket_data.get("year2") if bucket_data else None,
    )
    return WaterfallPayload(
        categories=list(categories),
        series_values=[(name, list(values)) for name, values in series_values],
        base_indices=base_indices,
        base_values=base_values,
        gathered_label_values={
            key: list(values) for key, values in gathered_label_values.items()
        },
    )


def compute_waterfall_payloads_for_all_labels(
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None,
    bucket_data: dict | None,
    template_chart,
    target_labels: list[str] | None = None,
) -> dict[str, WaterfallPayload]:
    labels = _normalize_target_level_labels(target_labels)
    if not labels:
        labels = _target_level_labels_from_gathered_df_with_filters(
            gathered_df,
            year1=bucket_data.get("year1") if bucket_data else None,
            year2=bucket_data.get("year2") if bucket_data else None,
            target_labels=bucket_data.get("target_labels") if bucket_data else None,
        )
    payloads_by_label = {}
    logger.info("Precomputing waterfall payloads for %d label(s).", len(labels))
    for label in labels:
        payload = compute_payload_for_label(
            gathered_df,
            scope_df,
            label,
            bucket_data,
            template_chart,
        )
        payloads_by_label[label] = payload
        logger.info(
            "Computed waterfall payload for %r: %d categories, checksum %.2f",
            label,
            len(payload.categories),
            _payload_checksum(payload.series_values),
        )
        out = Path("debug_waterfall_payloads.json")
        with out.open("w", encoding="utf-8") as f:
            json.dump(payloads_by_label, f, indent=2, ensure_ascii=False)

        print(f"[waterfall] wrote payload debug to: {out.resolve()}")
    return payloads_by_label


def _add_waterfall_chart_from_template(
    slide,
    template_slide,
    scope_df: pd.DataFrame | None,
    gathered_df: pd.DataFrame | None,
    target_level_label: str | None,
    bucket_data: dict | None,
    chart_name: str,
):
    template_shape = _waterfall_chart_shape_from_slide(template_slide, chart_name)
    if template_shape is None:
        raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
    template_chart = template_shape.chart
    template_series_names = [series.name for series in template_chart.series]
    label_columns = _capture_label_columns(
        _load_chart_workbook(template_chart).active,
        template_series_names,
    )
    (
        cd,
        categories,
        base_indices,
        base_values,
        _,
        gathered_label_values,
    ) = _build_waterfall_chart_data(
        template_chart,
        scope_df,
        gathered_df,
        target_level_label,
        bucket_data.get("labels") if bucket_data else None,
        bucket_data.get("values") if bucket_data else None,
        year1=bucket_data.get("year1") if bucket_data else None,
        year2=bucket_data.get("year2") if bucket_data else None,
    )
    chart_type = getattr(
        template_chart,
        "chart_type",
        getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED),
    )
    chart_shape = slide.shapes.add_chart(
        chart_type,
        template_shape.left,
        template_shape.top,
        template_shape.width,
        template_shape.height,
        cd,
    )
    chart_shape.name = getattr(template_shape, "name", chart_name)
    updated_wb = _load_chart_workbook(chart_shape.chart)
    total_rows = len(categories)
    _update_lab_base_label(
        label_columns,
        base_indices,
        base_values,
        total_rows,
    )
    _apply_label_columns(updated_wb.active, label_columns, total_rows)
    _ensure_negatives_column_positive(updated_wb.active)
    applied_headers = _apply_gathered_waterfall_labels(
        updated_wb.active,
        gathered_label_values,
        total_rows,
    )
    _update_all_waterfall_labs(
        updated_wb.active,
        base_indices,
        base_values,
        skip_headers=applied_headers,
    )
    _save_chart_workbook(chart_shape.chart, updated_wb)
    _update_waterfall_chart_caches(chart_shape.chart, updated_wb, categories)
    _update_waterfall_yoy_arrows(slide, base_values)
    return chart_shape


def _waterfall_template_marker(index: int) -> str:
    if index < 0:
        raise ValueError("Waterfall template index must be non-negative.")
    if index == 0:
        return "<Waterfall Template>"
    return f"<Waterfall Template{index + 1}>"


def _available_waterfall_template_slides(prs) -> list[tuple[str, object]]:
    slides = []
    idx = 0
    while True:
        marker = _waterfall_template_marker(idx)
        slide = _find_slide_by_marker(prs, marker)
        if slide is None:
            break
        slides.append((marker, slide))
        idx += 1
    return slides


def _normalize_target_level_labels(labels: list[str] | None) -> list[str]:
    unique_labels = []
    seen = set()
    for label in labels or []:
        if label is None:
            continue
        value = str(label).strip()
        if not value or value in seen:
            continue
        seen.add(value)
        unique_labels.append(value)
    return unique_labels


def _set_waterfall_slide_header(slide, label: str, marker_text: str | None = None) -> None:
    title_text = label
    replaced = False
    if marker_text:
        replaced = replace_text_in_slide_preserve_formatting(slide, marker_text, title_text)
        marker_plain = marker_text.strip("<>")
        replaced = (
            replace_text_in_slide_preserve_formatting(slide, marker_plain, title_text)
            or replaced
        )
    replaced = replace_text_in_slide_preserve_formatting(
        slide, "<Waterfall Template>", title_text
    ) or replaced
    replaced = (
        replace_text_in_slide_preserve_formatting(slide, "Waterfall Template", title_text)
        or replaced
    )
    if replaced:
        return
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.text = title_text
            return


def _update_waterfall_chart(
    slide,
    payload: WaterfallPayload,
) -> None:
    chart_shapes = [shape for shape in slide.shapes if shape.has_chart]
    if not chart_shapes:
        raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
    for chart_shape in chart_shapes:
        chart = chart_shape.chart
        series_names = [series.name for series in chart.series]
        label_columns = _capture_label_columns(_load_chart_workbook(chart).active, series_names)
        cd = _chart_data_from_payload(payload)
        chart.replace_data(cd)
        updated_wb = _load_chart_workbook(chart)
        total_rows = len(payload.categories)
        _update_lab_base_label(
            label_columns,
            payload.base_indices,
            payload.base_values,
            total_rows,
        )
        _apply_label_columns(updated_wb.active, label_columns, total_rows)
        _ensure_negatives_column_positive(updated_wb.active)
        applied_headers = _apply_gathered_waterfall_labels(
            updated_wb.active,
            payload.gathered_label_values,
            total_rows,
        )
        _update_all_waterfall_labs(
            updated_wb.active,
            payload.base_indices,
            payload.base_values,
            skip_headers=applied_headers,
        )
        _save_chart_workbook(chart, updated_wb)
        chart_type = getattr(
            chart,
            "chart_type",
            getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED),
        )
        if chart_type == getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED):
            _update_waterfall_chart_caches(chart, updated_wb, payload.categories)
        else:
            _update_chart_label_caches(chart, updated_wb)
    _update_waterfall_yoy_arrows(slide, payload.base_values)


def _resolve_target_level_label_value(
    gathered_df: pd.DataFrame | None,
    waterfall_targets: list[str] | None,
    bucket_data: dict | None,
) -> str | None:
    selected = [label for label in (waterfall_targets or []) if label and str(label).strip()]
    if selected:
        if len(selected) > 1:
            joined = ", ".join(str(label).strip() for label in selected)
            logger.info(
                "Multiple Target Level Label values selected; using joined string: %s",
                joined,
            )
            return joined
        return str(selected[0]).strip()
    if gathered_df is None or gathered_df.empty:
        return None
    year1 = bucket_data.get("year1") if bucket_data else None
    year2 = bucket_data.get("year2") if bucket_data else None
    target_labels = bucket_data.get("target_labels") if bucket_data else None
    labels = _target_level_labels_from_gathered_df_with_filters(
        gathered_df,
        year1=year1,
        year2=year2,
        target_labels=target_labels,
    )
    if not labels:
        return None
    if len(labels) > 1:
        joined = ", ".join(labels)
        logger.info(
            "Multiple Target Level Label values derived from gatheredCN10; using joined string: %s",
            joined,
        )
        return joined
    return labels[0]


def populate_category_waterfall(
    prs,
    gathered_df: pd.DataFrame,
    scope_df: pd.DataFrame | None = None,
    target_labels: list[str] | None = None,
    bucket_data: dict | None = None,
    modelled_in_value: str | None = None,
    metric_value: str | None = None,
):
    labels = _normalize_target_level_labels(target_labels)
    if not labels:
        labels = _target_level_labels_from_gathered_df_with_filters(
            gathered_df,
            year1=bucket_data.get("year1") if bucket_data else None,
            year2=bucket_data.get("year2") if bucket_data else None,
            target_labels=bucket_data.get("target_labels") if bucket_data else None,
        )
    if not labels:
        return

    available_slides = _available_waterfall_template_slides(prs)
    available_count = len(available_slides)
    if available_count == 0:
        raise ValueError("Could not find the <Waterfall Template> slide in the template.")
    if len(labels) > available_count:
        raise ValueError(
            "Need {needed} waterfall slides but only found {available} "
            "(<Waterfall Template>...<Waterfall Template{available}>) in template. "
            "Please add more duplicated template slides or use a larger template deck.".format(
                needed=len(labels),
                available=available_count,
            )
        )

    template_chart = _waterfall_chart_from_slide(available_slides[0][1], "Waterfall Template")
    if template_chart is None:
        raise ValueError("Could not find the waterfall chart on the <Waterfall Template> slide.")
    payloads_by_label = compute_waterfall_payloads_for_all_labels(
        gathered_df,
        scope_df,
        bucket_data,
        template_chart,
        target_labels=labels,
    )

    seen_partnames: set[str] = set()
    remaining_labels = labels.copy()
    for idx in range(len(labels)):
        marker_text, slide = available_slides[idx]
        resolved_label = _resolve_target_level_label_for_slide(slide, remaining_labels)
        if resolved_label is None:
            if not remaining_labels:
                raise ValueError("No remaining Target Level Label values to assign.")
            resolved_label = remaining_labels[0]
            logger.info(
                "No slide title/name found for %s; using ordered Target Level Label %r.",
                marker_text,
                resolved_label,
            )
        if resolved_label in remaining_labels:
            remaining_labels.remove(resolved_label)
        if resolved_label not in payloads_by_label:
            raise ValueError(
                f"Missing precomputed waterfall payload for Target Level Label {resolved_label!r}."
            )
        _ensure_unique_chart_parts_on_slide(slide, seen_partnames)
        _update_waterfall_axis_placeholders(
            prs,
            slide,
            target_level_label_value=resolved_label,
            modelled_in_value=modelled_in_value,
            metric_value=metric_value,
        )
        _update_waterfall_chart(
            slide,
            payloads_by_label[resolved_label],
        )
        _set_waterfall_slide_header(slide, resolved_label, marker_text=marker_text)

def build_pptx_from_template(
    template_bytes,
    df,
    target_brand=None,
    project_name=None,
    scope_df=None,
    product_description_df=None,
    waterfall_targets=None,
    bucket_data=None,
    modelled_in_value: str | None = None,
    metric_value: str | None = None,
):
    global _num_cache_warning_count
    _num_cache_warning_count = 0
    prs = Presentation(io.BytesIO(template_bytes))
    # Assume Slide 1 has TitleBox & SubTitle
    slide1 = prs.slides[0]
    set_text_by_name(slide1, "TitleBox", "Monthly Performance Summary")
    set_text_by_name(slide1, "SubTitle", "Auto-generated via Dash + python-pptx")
    if target_brand:
        replace_text_in_slide(slide1, "Target Brand", target_brand)
    if project_name == "MMx" and scope_df is not None:
        try:
            year_range = _format_study_year_range(scope_df)
        except Exception:
            year_range = None
        if year_range:
            replace_text_in_slide_preserve_formatting(slide1, "<RANGE>", year_range)
        generation_date = date.today().strftime("%b %d, %Y")
        replace_text_in_slide_preserve_formatting(slide1, "<DATE>", generation_date)
    remove_empty_placeholders(slide1)

    # Assume Slide 2 is for a KPI table and a chart
    slide2 = prs.slides[1] if len(prs.slides) > 1 else prs.slides.add_slide(prs.slide_layouts[5])

    # Simple KPIs (example): top 5 brands by value
    if "Brand" in df.columns and "Value" in df.columns:
        kpis = (
            df.groupby("Brand", as_index=False)["Value"].sum()
              .sort_values("Value", ascending=False)
              .head(5)
        )
        add_table(slide2, "Table_Summary", kpis)

        # Chart: share by Brand (editable)
        categories = kpis["Brand"].tolist()
        series = {"Value": kpis["Value"].tolist()}
        update_or_add_column_chart(slide2, "Chart_ShareByBrand", categories, series)

    remove_empty_placeholders(slide2)

    if project_name == "MMx" and len(prs.slides) > 3:
        slide4 = prs.slides[3]
        if scope_df is not None:
            try:
                modelled_category = modelled_category_from_scope_df(scope_df)
            except Exception:
                modelled_category = None
            if modelled_category:
                append_text_after_label(slide4, "Modelled Category:", modelled_category)
        if product_description_df is not None:
            try:
                target_dimensions = target_dimensions_from_product_description(
                    product_description_df
                )
            except Exception:
                target_dimensions = []
            if target_dimensions:
                append_paragraphs_after_label(
                    slide4,
                    "Modelled Category:",
                    target_dimensions,
                )
        if scope_df is not None:
            try:
                time_period, week_count = _format_modelling_period(df, scope_df)
            except Exception:
                time_period = None
                week_count = None
            if time_period and week_count is not None:
                set_time_period_text(slide4, "TIME PERIOD", time_period, week_count)
        remove_empty_placeholders(slide4)

    if project_name == "MMx":
        try:
            populate_category_waterfall(
                prs,
                df,
                scope_df,
                waterfall_targets,
                bucket_data,
                modelled_in_value,
                metric_value,
            )
        except Exception:
            logger.exception("Failed to populate category waterfall slides.")
            raise

    # Return bytes
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()

def render_upload_status(filename, success_label):
    if not filename:
        return html.Div("No file uploaded yet.", style={"color": "#888", "fontSize": "0.9rem"})

    return html.Div(
        [
            html.Div(
                style={
                    "height": "10px",
                    "backgroundColor": "#E5E7EB",
                    "borderRadius": "999px",
                    "overflow": "hidden",
                    "marginTop": "8px",
                },
                children=[
                    html.Div(
                        style={
                            "width": "100%",
                            "height": "100%",
                            "backgroundColor": "#22C55E",
                            "transition": "width 0.3s ease",
                        }
                    )
                ],
            ),
            html.Div(
                f"{success_label}: {filename}",
                style={"color": "#15803D", "fontSize": "0.9rem", "marginTop": "6px"},
            ),
        ]
    )

app.layout = html.Div(
    style={"maxWidth":"900px","margin":"40px auto","fontFamily":"Inter, system-ui"},
    children=[
        html.H2("PowerPoint Deck Automator (Dash + python-pptx)"),
        html.P("Upload your data, pick the project, and we will fill the matching PPTX template."),
        dcc.Store(id="bucket-metadata"),
        dcc.Store(id="bucket-config"),
        dcc.Store(id="bucket-deltas"),
        html.Div(
            [
                html.Label("Which project are you working on?"),
                dcc.Dropdown(
                    id="project-select",
                    options=[{"label": key, "value": key} for key in PROJECT_TEMPLATES],
                    placeholder="Select a project",
                    clearable=False,
                ),
            ],
            style={"marginBottom": "18px"},
        ),
        html.Div([
            html.Label("Upload gatheredCN10 file (.xlsx, .xlsb, or .csv):"),
            dcc.Upload(
                id="data-upload",
                children=html.Div(["Drag & Drop or ", html.A("Select File")]),
                accept=".xlsx,.xls,.xlsb,.csv",
                multiple=False,
                style={
                    "padding":"20px",
                    "border":"1px dashed #888",
                    "borderRadius":"12px",
                    "marginBottom":"6px",
                },
            ),
            html.Div(
                id="data-upload-status",
                children=render_upload_status(None, "gatheredCN10 upload complete"),
                style={"marginBottom":"12px"},
            ),
        ], style={"marginBottom":"18px"}),
        html.Div([
            html.Label("Upload scope file (.xlsx or .xlsb):"),
            dcc.Upload(
                id="scope-upload",
                children=html.Div(["Drag & Drop or ", html.A("Select File")]),
                accept=".xlsx,.xlsb",
                multiple=False,
                style={
                    "padding":"20px",
                    "border":"1px dashed #888",
                    "borderRadius":"12px",
                    "marginBottom":"6px",
                },
            ),
            html.Div(
                id="scope-upload-status",
                children=render_upload_status(None, "Scope upload complete"),
                style={"marginBottom":"12px"},
            ),
        ], style={"marginBottom":"18px"}),
        html.Div(
            [
                html.Label(
                    "Which Target Level Label values should be included in the waterfalls?"
                ),
                dcc.Checklist(
                    id="waterfall-targets",
                    options=[],
                    value=[],
                    labelStyle={"display": "block", "marginBottom": "6px"},
                    inputStyle={"marginRight": "8px"},
                ),
                html.Div(
                    id="waterfall-targets-status",
                    style={"color": "#6B7280", "fontSize": "0.9rem", "marginTop": "6px"},
                ),
            ],
            style={"marginBottom": "18px"},
        ),
        html.Div(
            [
                html.H3("Bucketed Waterfall Inputs"),
                html.Div(
                    [
                        html.Label("Year 1"),
                        dcc.Dropdown(
                            id="bucket-year1",
                            options=[],
                            placeholder="Select Year 1",
                            clearable=False,
                        ),
                    ],
                    style={"marginBottom": "12px"},
                ),
                html.Div(
                    [
                        html.Label("Year 2"),
                        dcc.Dropdown(
                            id="bucket-year2",
                            options=[],
                            placeholder="Select Year 2",
                            clearable=False,
                        ),
                    ],
                    style={"marginBottom": "12px"},
                ),
                html.Div(
                    [
                        html.Button(
                            "Select all buckets",
                            id="bucket-select-all",
                            n_clicks=0,
                            style={"padding": "6px 12px", "borderRadius": "8px"},
                        ),
                        html.Button(
                            "Clear all buckets",
                            id="bucket-clear-all",
                            n_clicks=0,
                            style={"padding": "6px 12px", "borderRadius": "8px"},
                        ),
                    ],
                    style={
                        "display": "flex",
                        "gap": "8px",
                        "marginBottom": "12px",
                        "flexWrap": "wrap",
                    },
                ),
                html.Div(id="bucket-group-controls"),
                html.Button(
                    "Apply Buckets to Waterfall",
                    id="apply-buckets",
                    n_clicks=0,
                    style={"padding": "10px 16px", "borderRadius": "10px"},
                ),
                html.Div(
                    id="bucket-status",
                    style={"color": "#6B7280", "fontSize": "0.9rem", "marginTop": "8px"},
                ),
            ],
            style={
                "marginBottom": "18px",
                "padding": "12px",
                "border": "1px solid #E5E7EB",
                "borderRadius": "12px",
            },
        ),

        html.Button("Generate Deck", id="go", n_clicks=0, style={"padding":"10px 16px","borderRadius":"10px"}),
        html.Div(id="status", style={"marginTop":"10px", "color":"#888"}),
        dcc.Download(id="download"),
    ]
)

@callback(
    Output("data-upload-status", "children"),
    Input("data-upload", "contents"),
    State("data-upload", "filename"),
)
def show_data_upload_status(contents, filename):
    if not contents:
        return render_upload_status(None, "gatheredCN10 upload complete")
    return render_upload_status(filename, "gatheredCN10 upload complete")

@callback(
    Output("scope-upload-status", "children"),
    Input("scope-upload", "contents"),
    State("scope-upload", "filename"),
)
def show_scope_upload_status(contents, filename):
    if not contents:
        return render_upload_status(None, "Scope upload complete")
    return render_upload_status(filename, "Scope upload complete")

@callback(
    Output("waterfall-targets", "options"),
    Output("waterfall-targets", "value"),
    Output("waterfall-targets-status", "children"),
    Input("data-upload", "contents"),
    State("data-upload", "filename"),
)
def populate_waterfall_targets(contents, filename):
    if not contents:
        return [], [], "Upload a gatheredCN10 file to load Target Level Label values."
    try:
        gathered_df = df_from_contents(contents, filename)
    except Exception as exc:
        return [], [], f"Error reading gatheredCN10 file: {exc}"
    if gathered_df is None or gathered_df.empty:
        return [], [], "The gatheredCN10 file is empty."
    try:
        target_labels = _target_level_labels_from_gathered_df(gathered_df)
    except Exception as exc:
        return [], [], f"Error finding Target Level Label values: {exc}"
    if not target_labels:
        return [], [], "No Target Level Label values were found in the gatheredCN10 file."
    options = [{"label": label, "value": label} for label in target_labels]
    return options, target_labels, f"Found {len(target_labels)} Target Level Label value(s)."


@callback(
    Output("bucket-metadata", "data"),
    Output("bucket-year1", "options"),
    Output("bucket-year1", "value"),
    Output("bucket-year2", "options"),
    Output("bucket-year2", "value"),
    Output("bucket-group-controls", "children"),
    Output("bucket-status", "children"),
    Input("data-upload", "contents"),
    State("data-upload", "filename"),
)
def populate_bucket_controls(contents, filename):
    if not contents:
        return (
            {},
            [],
            None,
            [],
            None,
            [],
            "Upload a gatheredCN10 file to configure bucket inputs.",
        )
    try:
        raw_df = raw_df_from_contents(contents, filename)
        data_df, metadata = _parse_two_row_header_dataframe(raw_df)
    except Exception as exc:
        return (
            {},
            [],
            None,
            [],
            None,
            [],
            f"Error parsing two-row headers: {exc}",
        )

    if not metadata.get("target_label_id"):
        return (
            metadata,
            [],
            None,
            [],
            None,
            [],
            "The gatheredCN10 file is missing the Target Label column.",
        )
    if not metadata.get("year_id"):
        return (
            metadata,
            [],
            None,
            [],
            None,
            [],
            "The gatheredCN10 file is missing the Year column.",
        )

    year_values = _unique_column_values(data_df, metadata["year_id"])
    year_options = [{"label": value, "value": value} for value in year_values]
    year1_default = year_values[0] if year_values else None
    year2_default = year_values[1] if len(year_values) > 1 else year1_default

    column_labels: dict[str, str] = {}
    column_groups: dict[str, str] = {}
    group_controls = []
    for group in metadata.get("group_order", []):
        columns = metadata.get("groups", {}).get(group, [])
        if not columns:
            continue
        label_counts: dict[str, int] = {}
        column_rows = []
        for column in columns:
            subheader = column.get("subheader") or "Unnamed"
            label_counts[subheader] = label_counts.get(subheader, 0) + 1
            label = subheader
            if label_counts[subheader] > 1:
                label = f"{subheader} ({label_counts[subheader]})"
            column_id = column["id"]
            column_labels[column_id] = label
            column_groups[column_id] = group
            column_rows.append(
                html.Div(
                    [
                        dcc.Checklist(
                            id={"type": "bucket-column", "column": column_id},
                            options=[{"label": label, "value": column_id}],
                            value=[column_id],
                            labelStyle={"display": "block", "marginBottom": "4px"},
                            inputStyle={"marginRight": "6px"},
                        ),
                    ],
                    style={
                        "display": "flex",
                        "alignItems": "center",
                        "justifyContent": "space-between",
                        "gap": "12px",
                        "marginBottom": "6px",
                    },
                )
            )
        group_controls.append(
            html.Div(
                [
                    html.Label(group, style={"fontWeight": "600"}),
                    html.Label(
                        "Target Label filter",
                        style={"fontSize": "0.85rem", "color": "#6B7280"},
                    ),
                    dcc.Checklist(
                        id={"type": "bucket-group-type", "group": group},
                        options=[
                            {"label": DISPLAY_LABEL["Own"], "value": "Own"},
                            {"label": DISPLAY_LABEL["Cross"], "value": "Cross"},
                        ],
                        value=[],
                        labelStyle={"display": "block", "marginBottom": "2px"},
                        inputStyle={"marginRight": "6px"},
                        style={"minWidth": "120px", "marginBottom": "8px"},
                    ),
                    *column_rows,
                ],
                style={
                    "marginBottom": "12px",
                    "padding": "8px",
                    "border": "1px solid #E5E7EB",
                    "borderRadius": "8px",
                },
            )
        )

    status = (
        f"Loaded {len(metadata.get('group_order', []))} bucket group(s)."
        if metadata.get("group_order")
        else "No bucket groups were found in the first header row."
    )
    return (
        {**metadata, "column_labels": column_labels, "column_groups": column_groups},
        year_options,
        year1_default,
        year_options,
        year2_default,
        group_controls,
        status,
    )


@callback(
    Output({"type": "bucket-group-type", "group": ALL}, "value"),
    Input("bucket-select-all", "n_clicks"),
    Input("bucket-clear-all", "n_clicks"),
    State({"type": "bucket-group-type", "group": ALL}, "id"),
    prevent_initial_call=True,
)
def update_bucket_group_types(select_all_clicks, clear_all_clicks, bucket_type_ids):
    if not bucket_type_ids:
        return []
    triggered = callback_context.triggered[0]["prop_id"].split(".")[0]
    if triggered == "bucket-clear-all":
        return [[] for _ in bucket_type_ids]
    if triggered == "bucket-select-all":
        return [["Own", "Cross"] for _ in bucket_type_ids]
    return no_update


@callback(
    Output("bucket-config", "data"),
    Output("bucket-deltas", "data"),
    Output("bucket-status", "children", allow_duplicate=True),
    Input("apply-buckets", "n_clicks"),
    State("data-upload", "contents"),
    State("data-upload", "filename"),
    State("bucket-year1", "value"),
    State("bucket-year2", "value"),
    State("bucket-metadata", "data"),
    State({"type": "bucket-column", "column": ALL}, "value"),
    State({"type": "bucket-column", "column": ALL}, "id"),
    State({"type": "bucket-group-type", "group": ALL}, "value"),
    State({"type": "bucket-group-type", "group": ALL}, "id"),
    prevent_initial_call=True,
)
def apply_bucket_selection(
    n_clicks,
    contents,
    filename,
    year1,
    year2,
    metadata,
    selections,
    selection_ids,
    bucket_types,
    bucket_type_ids,
):
    if not contents:
        return no_update, no_update, "Upload a gatheredCN10 file before applying buckets."
    if not metadata:
        return (
            no_update,
            no_update,
            "Bucket metadata is unavailable. Re-upload the gatheredCN10 file.",
        )
    if not year1 or not year2:
        return no_update, no_update, "Select Year 1 and Year 2 values before applying."

    try:
        raw_df = raw_df_from_contents(contents, filename)
        data_df, parsed_meta = _parse_two_row_header_dataframe(raw_df)
    except Exception as exc:
        return no_update, no_update, f"Error parsing gatheredCN10: {exc}"

    bucket_type_map: dict[str, list[str]] = {}
    for bucket_type, type_id in zip(bucket_types, bucket_type_ids):
        group = type_id.get("group")
        if group:
            bucket_type_map[group] = [value for value in bucket_type or [] if value]

    selected_columns = []
    for selection, selection_id in zip(selections, selection_ids):
        column_id = selection_id.get("column")
        if column_id and selection:
            selected_columns.append(column_id)

    column_groups = metadata.get("column_groups", {}) if metadata else {}
    group_columns: dict[str, list[str]] = {}
    for column_id in selected_columns:
        group = column_groups.get(column_id)
        if group:
            group_columns.setdefault(group, []).append(column_id)

    bucket_config: dict[str, dict[str, list[str]]] = {}
    for group in metadata.get("group_order", []):
        selected_group_columns = group_columns.get(group, [])
        target_labels = bucket_type_map.get(group)
        if target_labels is None:
            target_labels = []
        bucket_config[group] = {
            "target_labels": target_labels,
            "subheaders_included": selected_group_columns,
        }

    try:
        deltas = _compute_bucket_deltas(
            data_df,
            parsed_meta,
            bucket_config,
            year1,
            year2,
        )
    except Exception as exc:
        return no_update, no_update, f"Error computing bucket deltas: {exc}"

    labels = [group for group, _ in deltas]
    values = [value for _, value in deltas]
    selected_target_labels = []
    seen_target_labels = set()
    for config in bucket_config.values():
        for label in config.get("target_labels", []):
            if not label or label in seen_target_labels:
                continue
            seen_target_labels.add(label)
            selected_target_labels.append(label)

    bucket_data = {
        "labels": labels,
        "values": values,
        "year1": year1,
        "year2": year2,
        "target_labels": selected_target_labels,
    }
    return bucket_config, bucket_data, f"Applied {len(values)} bucket delta(s) to the waterfall."


@callback(
    Output("download","data"),
    Output("status","children"),
    Input("go","n_clicks"),
    State("data-upload","contents"),
    State("data-upload","filename"),
    State("scope-upload", "contents"),
    State("scope-upload", "filename"),
    State("project-select", "value"),
    State("waterfall-targets", "value"),
    State("bucket-deltas", "data"),
    prevent_initial_call=True
)
def generate_deck(
    n_clicks,
    data_contents,
    data_name,
    scope_contents,
    scope_name,
    project_name,
    waterfall_targets,
    bucket_data,
):
    if not data_contents or not project_name:
        return no_update, "Please upload the data file and select a project."

    template_path = PROJECT_TEMPLATES.get(project_name)
    if not template_path or not template_path.exists():
        return no_update, "The selected project template could not be found."
    try:
        df = df_from_contents(data_contents, data_name)
        scope_df = None
        product_description_df = None
        project_details_df = None
        modelled_in_value = None
        metric_value = None
        if scope_contents:
            try:
                scope_df = scope_df_from_contents(scope_contents, scope_name)
            except Exception:
                scope_df = None
            try:
                product_description_df = product_description_df_from_contents(
                    scope_contents, scope_name
                )
            except Exception:
                product_description_df = None
            try:
                project_details_df = project_details_df_from_contents(
                    scope_contents, scope_name
                )
            except Exception:
                project_details_df = None
        if project_details_df is not None:
            modelled_in_value = _project_detail_value_from_df(
                project_details_df,
                "modelled in",
                [
                    "Sales will be modelled in",
                    "Sales will be modeled in",
                    "Sales modelled in",
                    "Sales modeled in",
                ],
                "Sales will be modelled in",
            )
            metric_value = _project_detail_value_from_df(
                project_details_df,
                "metric",
                [
                    "Volume metric (unique per dataset)",
                    "Volume metric unique per dataset",
                    "Volume metric",
                ],
                "Volume metric (unique per dataset)",
            )
        target_brand = target_brand_from_scope_df(scope_df)
        template_bytes = template_path.read_bytes()

        pptx_bytes = build_pptx_from_template(
            template_bytes,
            df,
            target_brand,
            project_name,
            scope_df,
            product_description_df,
            waterfall_targets,
            bucket_data,
            modelled_in_value,
            metric_value,
        )
        return dcc.send_bytes(lambda buff: buff.write(pptx_bytes), "deck.pptx"), "Building deck..."

    except Exception as exc:
        logger.exception("Deck generation failed.")
        message = str(exc).strip()
        if not message:
            message = "Unknown error. Check server logs for details."
        return no_update, f"Error ({type(exc).__name__}): {message}"

# Important: Dash's dcc.send_bytes expects a writer function; we provide inline:
def _writer(f):
    pass

# Patch: we pass a lambda that writes nothing (handled internally). To attach bytes, we can use:
# return dcc.send_bytes(lambda b: b.write(pptx_bytes), "deck.pptx")

# Fix the callback to use the writer properly:
@callback(
    Output("download","data", allow_duplicate=True),
    Input("status","children"),
    State("data-upload","contents"),
    prevent_initial_call=True
)
def finalize_download(status_text, data_contents):
    # This is a no-op; left for clarity in a larger app. In the minimal example above,
    # you can directly return the 'dcc.send_bytes' with the actual bytes.
    return no_update

if __name__ == "__main__":
    app.run(debug=True)





