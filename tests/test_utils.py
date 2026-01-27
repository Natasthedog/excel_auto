from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from app import build_pptx_from_template


def build_sample_dataframe(labels: list[str], include_brand: bool = True) -> pd.DataFrame:
    rows = []
    for idx, label in enumerate(labels):
        for year, actuals in (("Year1", 100 + idx * 10), ("Year2", 120 + idx * 10)):
            row = {
                "Target Level Label": label,
                "Target Label": "Own",
                "Year": year,
                "Actuals": actuals,
            }
            if include_brand:
                row["Brand"] = f"Brand {idx + 1}"
                row["Value"] = actuals
            rows.append(row)
    return pd.DataFrame(rows)


def write_excel(df: pd.DataFrame, path: Path) -> None:
    df.to_excel(path, index=False)


def _add_textbox(slide, name: str, text: str, top_offset: float) -> None:
    left = Inches(0.5)
    top = Inches(top_offset)
    width = Inches(9)
    height = Inches(0.4)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.name = name
    textbox.text = text


def _add_waterfall_chart(slide, left_offset: float = 0.0) -> None:
    chart_data = ChartData()
    chart_data.categories = ["Base", "Change", "Total"]
    chart_data.add_series("Base", (100, 0, 110))
    chart_data.add_series("Positives", (0, 10, 0))
    chart_data.add_series("Negatives", (0, -5, 0))
    chart_data.add_series("Blanks", (0, 100, 105))
    chart_type = getattr(XL_CHART_TYPE, "WATERFALL", XL_CHART_TYPE.COLUMN_STACKED)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(0.8 + left_offset),
        Inches(2.2),
        Inches(8.5),
        Inches(3.5),
        chart_data,
    )
    chart_shape.name = "Waterfall Template"
    chart = chart_shape.chart
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True
    plot.data_labels.show_category_name = True


def build_test_template(
    path: Path,
    waterfall_slide_count: int,
    waterfall_chart_count: int = 1,
) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide1 = prs.slides.add_slide(blank_layout)
    _add_textbox(slide1, "TitleBox", "Title", 0.5)
    _add_textbox(slide1, "SubTitle", "Subtitle", 1.1)

    prs.slides.add_slide(blank_layout)
    prs.slides.add_slide(blank_layout)
    prs.slides.add_slide(blank_layout)

    for idx in range(waterfall_slide_count):
        marker = "<Waterfall Template>" if idx == 0 else f"<Waterfall Template{idx + 1}>"
        slide = prs.slides.add_slide(blank_layout)
        _add_textbox(slide, f"WaterfallMarker{idx + 1}", marker, 0.4)
        _add_textbox(slide, f"WaterfallLabel{idx + 1}", "Label: <Target Level Label>", 0.9)
        _add_textbox(slide, f"WaterfallModelled{idx + 1}", "Modelled in: <modelled in>", 1.4)
        _add_textbox(slide, f"WaterfallMetric{idx + 1}", "Metric: <metric>", 1.9)
        _add_textbox(slide, f"WaterfallStyle{idx + 1}", f"Template Style {idx + 1}", 6.1)
        for chart_idx in range(waterfall_chart_count):
            _add_waterfall_chart(slide, left_offset=9.2 * chart_idx)

    prs.save(path)


def build_deck_bytes(
    template_path: Path,
    df: pd.DataFrame,
    waterfall_targets: list[str] | None,
    modelled_in_value: str = "US",
    metric_value: str = "Units",
) -> bytes:
    return build_pptx_from_template(
        template_path.read_bytes(),
        df,
        target_brand=None,
        project_name="MMx",
        scope_df=None,
        product_description_df=None,
        waterfall_targets=waterfall_targets,
        bucket_data=None,
        modelled_in_value=modelled_in_value,
        metric_value=metric_value,
    )
