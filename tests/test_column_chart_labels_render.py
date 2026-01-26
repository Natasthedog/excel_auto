from __future__ import annotations

import io
import zipfile
import xml.etree.ElementTree as ET

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from test_utils import build_deck_bytes, build_sample_dataframe


def _build_column_chart_template(path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    prs.slides.add_slide(blank_layout)
    slide2 = prs.slides.add_slide(blank_layout)

    chart_data = ChartData()
    chart_data.categories = ["Old A", "Old B", "Old C"]
    chart_data.add_series("Value", (10, 20, 30))
    chart_shape = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(2),
        Inches(8),
        Inches(4.5),
        chart_data,
    )
    chart_shape.name = "Chart_ShareByBrand"
    chart = chart_shape.chart
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_value = True

    prs.save(path)


def _insert_value_from_cells_labels(template_path) -> None:
    template_bytes = template_path.read_bytes()
    with zipfile.ZipFile(io.BytesIO(template_bytes)) as zf:
        chart_files = [
            name
            for name in zf.namelist()
            if name.startswith("ppt/charts/chart") and name.endswith(".xml")
        ]
        assert chart_files
        chart_name = chart_files[0]
        chart_xml = zf.read(chart_name)

    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    ns_c15 = "http://schemas.microsoft.com/office/drawing/2012/chart"
    ET.register_namespace("c", ns_c)
    ET.register_namespace("c15", ns_c15)
    root = ET.fromstring(chart_xml)
    ns = {"c": ns_c}
    series_node = root.find(".//c:ser", ns)
    assert series_node is not None
    d_lbls = series_node.find("c:dLbls", ns)
    if d_lbls is None:
        d_lbls = ET.SubElement(series_node, f"{{{ns_c}}}dLbls")

    tx = ET.SubElement(d_lbls, f"{{{ns_c}}}tx")
    str_ref = ET.SubElement(tx, f"{{{ns_c}}}strRef")
    f_node = ET.SubElement(str_ref, f"{{{ns_c}}}f")
    f_node.text = "Sheet1!$A$2:$A$4"
    str_cache = ET.SubElement(str_ref, f"{{{ns_c}}}strCache")
    pt_count = ET.SubElement(str_cache, f"{{{ns_c}}}ptCount")
    pt_count.set("val", "3")
    for idx, value in enumerate(["old-1", "old-2", "old-3"]):
        pt = ET.SubElement(str_cache, f"{{{ns_c}}}pt", idx=str(idx))
        v = ET.SubElement(pt, f"{{{ns_c}}}v")
        v.text = value

    c15_range = ET.SubElement(d_lbls, f"{{{ns_c15}}}datalabelsRange")
    c15_f = ET.SubElement(c15_range, f"{{{ns_c15}}}f")
    c15_f.text = "Sheet1!$A$2:$A$4"
    c15_cache = ET.SubElement(c15_range, f"{{{ns_c15}}}dlblRangeCache")
    c15_pt_count = ET.SubElement(c15_cache, f"{{{ns_c15}}}ptCount")
    c15_pt_count.set("val", "3")
    for idx, value in enumerate(["old-1", "old-2", "old-3"]):
        pt = ET.SubElement(c15_cache, f"{{{ns_c15}}}pt", idx=str(idx))
        v = ET.SubElement(pt, f"{{{ns_c15}}}v")
        v.text = value

    updated_chart_xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)

    output = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(template_bytes)) as zf:
        with zipfile.ZipFile(output, "w") as zf_out:
            for info in zf.infolist():
                data = zf.read(info.filename)
                if info.filename == chart_name:
                    data = updated_chart_xml
                zf_out.writestr(info, data)
    template_path.write_bytes(output.getvalue())


def test_column_chart_value_from_cells_labels_update(tmp_path) -> None:
    template_path = tmp_path / "template.pptx"
    _build_column_chart_template(template_path)
    _insert_value_from_cells_labels(template_path)

    df = build_sample_dataframe(["Alpha", "Beta", "Gamma"], include_brand=True)
    pptx_bytes = build_deck_bytes(template_path, df, waterfall_targets=[])

    kpis = (
        df.groupby("Brand", as_index=False)["Value"].sum()
        .sort_values("Value", ascending=False)
        .head(5)
    )
    expected_labels = kpis["Brand"].tolist()

    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as zf:
        chart_files = [
            name
            for name in zf.namelist()
            if name.startswith("ppt/charts/chart") and name.endswith(".xml")
        ]
        assert chart_files
        chart_xml = zf.read(chart_files[0])

    ns = {
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
        "c15": "http://schemas.microsoft.com/office/drawing/2012/chart",
    }
    root = ET.fromstring(chart_xml)
    series_node = root.find(".//c:ser", ns)
    assert series_node is not None
    str_cache = series_node.find(".//c:dLbls//c:tx//c:strRef/c:strCache", ns)
    assert str_cache is not None
    points = str_cache.findall("c:pt", ns)
    values = [
        (pt.find("c:v", ns).text or "") if pt.find("c:v", ns) is not None else ""
        for pt in points
    ]
    assert values == expected_labels

    c15_cache = series_node.find(".//c15:datalabelsRange/c15:dlblRangeCache", ns)
    assert c15_cache is not None
    c15_points = c15_cache.findall("c15:pt", ns)
    c15_values = [
        (pt.find("c15:v", ns).text or "") if pt.find("c15:v", ns) is not None else ""
        for pt in c15_points
    ]
    assert c15_values == expected_labels
