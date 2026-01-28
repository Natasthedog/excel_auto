from __future__ import annotations

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app import _resolve_target_level_label_for_slide, _waterfall_base_values


def _make_slide_with_title(title_text: str):
    prs = Presentation()
    title_layout = None
    for layout in prs.slide_layouts:
        if any(
            placeholder.placeholder_format.type == PP_PLACEHOLDER.TITLE
            for placeholder in layout.placeholders
        ):
            title_layout = layout
            break
    if title_layout is None:
        title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title_text
    return slide


def test_resolve_target_level_label_from_slide_title() -> None:
    slide = _make_slide_with_title("  beta ")
    resolved = _resolve_target_level_label_for_slide(slide, ["Alpha", "Beta"])
    assert resolved == "Beta"


def test_resolve_target_level_label_returns_none_for_unmatched_title() -> None:
    slide = _make_slide_with_title("<Waterfall Template>")
    resolved = _resolve_target_level_label_for_slide(slide, ["Alpha", "Beta"])
    assert resolved is None


def test_waterfall_base_values_use_own_target_label() -> None:
    df = pd.DataFrame(
        [
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year1", "Actuals": 10},
            {"Target Level Label": "Alpha", "Target Label": "Own", "Year": "Year2", "Actuals": 20},
            {
                "Target Level Label": "Alpha",
                "Target Label": "Cross",
                "Year": "Year1",
                "Actuals": 100,
            },
            {
                "Target Level Label": "Alpha",
                "Target Label": "Cross",
                "Year": "Year2",
                "Actuals": 150,
            },
        ]
    )
    year1_total, year2_total = _waterfall_base_values(
        df,
        "Alpha",
        year1="Year1",
        year2="Year2",
    )
    assert year1_total == 10
    assert year2_total == 20
